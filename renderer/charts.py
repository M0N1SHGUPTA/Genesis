"""
renderer/charts.py — Native python-pptx chart generation.

Charts are embedded as real Office chart objects — they stay fully editable
in PowerPoint, Google Slides, and LibreOffice. No images or screenshots.

Supported chart types: bar (clustered column), pie, line with markers, area.
"""

from __future__ import annotations

import logging
from dataclasses import dataclass

from pptx.chart.data import CategoryChartData   # builds chart data with categories + series
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

import config
from renderer.utils import add_textbox

logger = logging.getLogger(__name__)

# Map the blueprint's string chart type → python-pptx XL_CHART_TYPE enum value.
# XL_CHART_TYPE values correspond directly to Office chart type identifiers.
_CHART_TYPE_MAP = {
    "bar":  XL_CHART_TYPE.COLUMN_CLUSTERED,   # vertical bar chart (most common)
    "pie":  XL_CHART_TYPE.PIE,                # pie / donut chart
    "line": XL_CHART_TYPE.LINE_MARKERS,       # line chart with data point markers
    "area": XL_CHART_TYPE.AREA,               # filled area chart for cumulative trends
}


@dataclass
class ChartPosition:
    """Stores position and size for a chart shape in EMU units.

    Using a dataclass avoids passing 4 individual positional arguments everywhere.
    """
    left: int    # distance from slide left edge (EMU)
    top: int     # distance from slide top edge (EMU)
    width: int   # chart width (EMU)
    height: int  # chart height (EMU)


# ---------------------------------------------------------------------------
# Public entry point (called by engine.py)
# ---------------------------------------------------------------------------

def render_chart_slide(slide, slide_data: dict, slide_num: int) -> None:
    """Render a complete chart slide: title + chart + optional caption + slide number.

    If the chart cannot be rendered (bad data, unsupported type, etc.) a
    placeholder error message is shown instead — the slide is never left blank.

    Args:
        slide: python-pptx slide object (already added to the presentation).
        slide_data: Single slide dict from the blueprint JSON.
        slide_num: Slide number to display in the bottom-right corner.
    """
    from renderer.utils import add_slide_title, add_slide_number
    from pptx.enum.text import PP_ALIGN

    # --- Title ---
    title = slide_data.get("title", "")
    add_slide_title(slide, title)

    # --- Chart area sizing ---
    # If there's a caption, shrink the chart height to leave room for it below
    caption = slide_data.get("caption", "")
    caption_height = Inches(0.4) if caption else 0

    pos = ChartPosition(
        left=config.MARGIN_LEFT,
        top=config.CONTENT_TOP,
        width=config.CONTENT_WIDTH,
        # Reserve space for caption + gap at the bottom
        height=config.CONTENT_HEIGHT - caption_height - config.ELEMENT_GAP,
    )

    # --- Chart ---
    chart_type_str = slide_data.get("chart_type", "bar").lower()
    data = slide_data.get("data", {})

    try:
        _add_chart(slide, chart_type_str, data, pos)
    except Exception as exc:
        # Something went wrong (missing data, API error, etc.)
        # Show a visible error message instead of an empty slide
        logger.warning("Chart render failed (%s): %s", chart_type_str, exc)
        add_textbox(
            slide, f"[Chart could not be rendered: {exc}]",
            left=pos.left, top=pos.top,
            width=pos.width, height=pos.height,
            font_size=config.BODY_FONT_SIZE,
            color=config.COLOR_TEXT_MUTED,
        )

    # --- Optional caption below the chart ---
    if caption:
        add_textbox(
            slide, caption,
            left=config.MARGIN_LEFT,
            top=config.CONTENT_TOP + pos.height + config.ELEMENT_GAP,
            width=config.CONTENT_WIDTH,
            height=Inches(0.35),
            font_size=config.CAPTION_FONT_SIZE,
            color=config.COLOR_TEXT_MUTED,
            align=PP_ALIGN.CENTER,   # captions are centred under the chart
        )

    # --- Slide number ---
    add_slide_number(slide, slide_num)


# ---------------------------------------------------------------------------
# Internal chart builder
# ---------------------------------------------------------------------------

def _add_chart(slide, chart_type_str: str, data: dict, pos: ChartPosition) -> None:
    """Build and embed a native Office chart into the slide.

    Steps:
      1. Look up the XL_CHART_TYPE from the string name
      2. Build a CategoryChartData object with categories + series
      3. Call slide.shapes.add_chart() to embed the chart
      4. Apply theme-consistent styling

    Args:
        slide: python-pptx slide object.
        chart_type_str: One of "bar", "pie", "line", "area".
        data: Dict with "categories" (list) and "series" (list of {name, values}).
        pos: ChartPosition holding left/top/width/height in EMU.
    """
    # Map string → enum; default to clustered column if type is unknown
    xl_type = _CHART_TYPE_MAP.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

    categories = data.get("categories", [])
    series_list = data.get("series", [])

    if not categories or not series_list:
        raise ValueError("Chart data missing 'categories' or 'series'.")

    # CategoryChartData holds the data table that Office reads to draw the chart
    chart_data = CategoryChartData()
    chart_data.categories = [str(c) for c in categories]   # category axis labels

    for series in series_list:
        name = series.get("name", "Series")
        values = series.get("values", [])

        # Coerce all values to float — non-numeric values become 0 to avoid crashes
        safe_values = []
        for v in values:
            try:
                safe_values.append(float(v))
            except (TypeError, ValueError):
                safe_values.append(0.0)   # replace bad data with zero

        chart_data.add_series(name, safe_values)

    # add_chart() returns a GraphicFrame; .chart gives access to the Chart object
    graphic_frame = slide.shapes.add_chart(
        xl_type,
        pos.left, pos.top,
        pos.width, pos.height,
        chart_data,
    )
    chart = graphic_frame.chart

    # Apply visual styling to match the template theme
    _style_chart(chart, chart_type_str, series_list)


def _style_chart(chart, chart_type_str: str, series_list: list[dict]) -> None:
    """Apply theme-consistent visual styling to an embedded chart.

    Covers: legend visibility, series colors, axis fonts, background transparency.

    Args:
        chart: python-pptx Chart object (from graphic_frame.chart).
        chart_type_str: Used for type-specific tweaks (e.g. bar gap width).
        series_list: The raw series list from the blueprint (used to check count).
    """
    # Show the legend only when there are multiple data series
    # (a single series doesn't need a legend — the title already identifies it)
    chart.has_legend = len(series_list) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM   # put legend below the chart
        chart.legend.include_in_layout = False               # don't shrink chart area for legend

    # Apply theme colors to each series in order from config.CHART_COLORS
    try:
        plot = chart.plots[0]
        for i, series in enumerate(plot.series):
            # Cycle through chart colors if there are more series than colors
            color = config.CHART_COLORS[i % len(config.CHART_COLORS)]
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = color
    except Exception as exc:
        logger.debug("Could not apply series colors: %s", exc)

    # Bar-specific: tighten the gap between bars for a cleaner, modern look
    if chart_type_str == "bar":
        try:
            chart.plots[0].gap_width = 80   # default is 150; smaller = wider bars
        except Exception:
            pass

    # Set axis tick label font size and color (value axis = y-axis / numbers)
    try:
        chart.value_axis.tick_labels.font.size = Pt(10)
        chart.value_axis.tick_labels.font.color.rgb = config.COLOR_TEXT_MUTED
    except Exception:
        pass   # some chart types (e.g. pie) don't have a value axis

    # Category axis = x-axis / category labels
    try:
        chart.category_axis.tick_labels.font.size = Pt(10)
        chart.category_axis.tick_labels.font.color.rgb = config.COLOR_TEXT_MUTED
    except Exception:
        pass

    # Make chart area background transparent so the slide background shows through
    try:
        chart.chart_area.format.fill.background()
    except Exception:
        pass

    # Make plot area (the inner data region) background transparent too
    try:
        chart.plot_area.format.fill.background()
    except Exception:
        pass
