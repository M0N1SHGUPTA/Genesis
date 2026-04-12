"""
renderer/design_system.py - Target-driven visual design extraction.

This module treats the decks in /target as the source of truth. It matches the
selected template to the closest reference deck by theme + layout signature,
then extracts the typography, color roles, spacing, and layout frames needed by
the renderer.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from functools import lru_cache
from pathlib import Path
import re
from statistics import median

from lxml import etree
from pptx import Presentation

EMU_PER_INCH = 914400
NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def _to_inches(value: int) -> float:
    return round(value / EMU_PER_INCH, 3)


def _as_int(value: float) -> int:
    return int(round(value))


@dataclass(frozen=True)
class TextStyle:
    name: str
    size: float
    bold: bool = False
    italic: bool = False
    color: str | None = None
    align: str = "left"
    line_spacing: float = 1.1
    paragraph_after_pt: float = 0.0


@dataclass(frozen=True)
class BoxSpec:
    left: float
    top: float
    width: float
    height: float


@dataclass(frozen=True)
class FooterSpec:
    enabled: bool
    left_mode: str | None = None
    left_box: BoxSpec | None = None
    number_box: BoxSpec | None = None
    text_style: TextStyle | None = None
    number_style: TextStyle | None = None


@dataclass(frozen=True)
class DesignSystem:
    family_id: str
    reference_path: str
    template_signature: str
    colors: dict[str, str]
    fonts: dict[str, TextStyle]
    spacing: dict[str, float]
    layouts: dict[str, BoxSpec]
    footer: FooterSpec
    metadata: dict[str, str] = field(default_factory=dict)


@dataclass(frozen=True)
class _DeckSignature:
    major_font: str
    minor_font: str
    layout_signature: str

    @property
    def key(self) -> str:
        return f"{self.major_font}|{self.minor_font}|{self.layout_signature}"


def _theme_info(prs: Presentation) -> tuple[str, str, dict[str, str]]:
    master = prs.slide_masters[0]
    theme_part = master.part.part_related_by(
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
    )
    root = etree.fromstring(theme_part.blob)
    font_scheme = root.find(".//a:fontScheme", NS)
    major = font_scheme.find("./a:majorFont/a:latin", NS).get("typeface") if font_scheme is not None else "Calibri"
    minor = font_scheme.find("./a:minorFont/a:latin", NS).get("typeface") if font_scheme is not None else "Calibri"

    colors: dict[str, str] = {}
    clr_scheme = root.find(".//a:clrScheme", NS)
    if clr_scheme is not None:
        for child in clr_scheme:
            tag = child.tag.split("}")[-1]
            srgb = child.find("./a:srgbClr", NS)
            sys_clr = child.find("./a:sysClr", NS)
            if srgb is not None:
                colors[tag] = srgb.get("val", "")
            elif sys_clr is not None:
                colors[tag] = sys_clr.get("lastClr", "")
    return major, minor, colors


def _signature_for_presentation(prs: Presentation) -> _DeckSignature:
    major, minor, _ = _theme_info(prs)
    layouts = "|".join(layout.name for layout in prs.slide_layouts)
    return _DeckSignature(major_font=major, minor_font=minor, layout_signature=layouts)


def _text_shapes(slide) -> list[dict]:
    shapes: list[dict] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text.strip().replace("\n", " | ")
        if not text:
            continue

        best_run = None
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if not run.text.strip():
                    continue
                size = run.font.size.pt if run.font.size else None
                if best_run is None or (size or 0) > (best_run[2] or 0):
                    best_run = (
                        run.font.name,
                        paragraph.alignment,
                        size,
                        bool(run.font.bold),
                    )

        shapes.append(
            {
                "text": text,
                "left": _to_inches(shape.left),
                "top": _to_inches(shape.top),
                "width": _to_inches(shape.width),
                "height": _to_inches(shape.height),
                "font_name": best_run[0] if best_run else None,
                "align": str(best_run[1]) if best_run and best_run[1] is not None else None,
                "font_size": best_run[2] if best_run else None,
                "bold": best_run[3] if best_run else False,
            }
        )
    return sorted(shapes, key=lambda item: (item["top"], item["left"]))


def _median_box(boxes: list[dict], *, default: BoxSpec) -> BoxSpec:
    if not boxes:
        return default
    return BoxSpec(
        left=median(item["left"] for item in boxes),
        top=median(item["top"] for item in boxes),
        width=median(item["width"] for item in boxes),
        height=median(item["height"] for item in boxes),
    )


def _first_matching_text(slide_shapes: list[dict], predicate) -> dict | None:
    for shape in slide_shapes:
        if predicate(shape):
            return shape
    return None


def _numeric_footer_box(slides: list, default: BoxSpec) -> BoxSpec:
    footer_boxes: list[dict] = []
    for slide in slides:
        for shape in _text_shapes(slide):
            text = shape["text"].strip()
            if (
                text.isdigit()
                and 0 < len(text) <= 2
                and shape["top"] >= 7.0
                and shape["left"] >= 9.0
            ):
                footer_boxes.append(shape)
    return _median_box(footer_boxes, default=default)


def _footer_label_box(slides: list, default: BoxSpec) -> BoxSpec | None:
    label_boxes: list[dict] = []
    for slide in slides:
        for shape in _text_shapes(slide):
            text = shape["text"].strip()
            if (
                len(text) > 10
                and not text.isdigit()
                and shape["top"] >= 7.0
                and shape["left"] <= 1.0
                and shape["height"] <= 0.2
            ):
                label_boxes.append(shape)
    if not label_boxes:
        return None
    return _median_box(label_boxes, default=default)


def _family_defaults(family_id: str) -> dict:
    if family_id == "oranienbaum-manrope":
        return {
            "cover_title_size": 36.0,
            "cover_subtitle_size": 14.0,
            "title_size": 24.0,
            "body_size": 14.0,
            "caption_size": 10.5,
            "card_heading_size": 16.0,
            "number_size": 24.0,
            "section_number_size": 56.0,
            "footer_left_mode": None,
            "inner_padding": 0.18,
            "card_gap": 0.22,
            "title_to_body": 0.42,
            "card_bg": "FFF7F7",
            "muted": "666666",
        }
    if family_id == "libre-inter":
        return {
            "cover_title_size": 40.0,
            "cover_subtitle_size": 15.0,
            "title_size": 22.0,
            "body_size": 13.0,
            "caption_size": 10.0,
            "card_heading_size": 14.0,
            "number_size": 40.0,
            "section_number_size": 40.0,
            "footer_left_mode": "presentation_title",
            "inner_padding": 0.16,
            "card_gap": 0.18,
            "title_to_body": 0.35,
            "card_bg": "FFF8F8",
            "muted": "6B7280",
        }
    return {
        "cover_title_size": 72.0,
        "cover_subtitle_size": 22.0,
        "title_size": 16.0,
        "body_size": 11.5,
        "caption_size": 10.0,
        "card_heading_size": 16.0,
        "number_size": 20.0,
        "section_number_size": 48.0,
        "footer_left_mode": None,
        "inner_padding": 0.14,
        "card_gap": 0.18,
        "title_to_body": 0.28,
        "card_bg": "F4F8EC",
        "muted": "5A6A56",
    }


def _family_id(signature: _DeckSignature) -> str:
    major = signature.major_font.lower()
    minor = signature.minor_font.lower()
    if "oranienbaum" in major and "manrope" in minor:
        return "oranienbaum-manrope"
    if "libre baskerville" in major and "inter" in minor:
        return "libre-inter"
    return "cambria"


def _extract_design_system(reference_path: str) -> DesignSystem:
    prs = Presentation(reference_path)
    signature = _signature_for_presentation(prs)
    family_id = _family_id(signature)
    defaults = _family_defaults(family_id)
    major_font, minor_font, theme_colors = _theme_info(prs)

    slides = list(prs.slides)
    cover_shapes = _text_shapes(slides[0]) if slides else []
    cover_title = cover_shapes[0] if cover_shapes else None
    cover_subtitle = cover_shapes[1] if len(cover_shapes) > 1 else None

    title_shapes: list[dict] = []
    for slide in slides[1:-1]:
        candidate = _first_matching_text(
            _text_shapes(slide),
            lambda item: item["top"] <= 1.25 and len(item["text"]) > 6 and item["width"] >= 3.0,
        )
        if candidate is not None:
            title_shapes.append(candidate)

    chart_boxes = [
        {
            "left": _to_inches(shape.left),
            "top": _to_inches(shape.top),
            "width": _to_inches(shape.width),
            "height": _to_inches(shape.height),
        }
        for slide in slides
        for shape in slide.shapes
        if str(shape.shape_type) == "CHART (3)"
    ]
    table_boxes = [
        {
            "left": _to_inches(shape.left),
            "top": _to_inches(shape.top),
            "width": _to_inches(shape.width),
            "height": _to_inches(shape.height),
        }
        for slide in slides
        for shape in slide.shapes
        if str(shape.shape_type) == "TABLE (19)"
    ]

    title_box = _median_box(
        title_shapes,
        default=BoxSpec(left=0.4, top=0.5, width=12.0, height=0.8),
    )
    chart_box = _median_box(
        chart_boxes,
        default=BoxSpec(left=0.55, top=1.55, width=12.1, height=5.1),
    )
    table_box = _median_box(
        table_boxes,
        default=BoxSpec(left=0.5, top=1.1, width=12.2, height=5.5),
    )
    cover_title_box = _median_box(
        [cover_title] if cover_title else [],
        default=BoxSpec(left=0.4, top=3.1, width=8.8, height=0.8),
    )
    cover_subtitle_box = _median_box(
        [cover_subtitle] if cover_subtitle else [],
        default=BoxSpec(left=0.4, top=4.2, width=8.8, height=0.8),
    )
    footer_number_box = _numeric_footer_box(
        slides[1:-1],
        default=BoxSpec(left=9.94, top=7.22, width=3.0, height=0.13),
    )
    footer_left_box = _footer_label_box(
        slides[1:-1],
        default=BoxSpec(left=0.38, top=7.22, width=3.4, height=0.14),
    )

    primary = theme_colors.get("accent1") or "EF4444"
    secondary = theme_colors.get("dk2") or theme_colors.get("accent2") or "2C2C2C"
    background = theme_colors.get("lt1") or "FFFFFF"
    light = theme_colors.get("lt2") or theme_colors.get("accent2") or defaults["card_bg"]

    content_left = title_box.left
    content_top = round(title_box.top + title_box.height + defaults["title_to_body"], 3)
    margin_bottom = round(7.5 - (footer_number_box.top + footer_number_box.height) + 0.03, 3)

    fonts = {
        "cover_title": TextStyle(name=major_font, size=defaults["cover_title_size"], bold=True, color=primary),
        "cover_subtitle": TextStyle(name=minor_font, size=defaults["cover_subtitle_size"], color=secondary, align="left"),
        "title": TextStyle(name=major_font, size=defaults["title_size"], bold=True, color=secondary),
        "section_number": TextStyle(name=major_font, size=defaults["section_number_size"], bold=True, color=primary),
        "body": TextStyle(name=minor_font, size=defaults["body_size"], color=secondary, paragraph_after_pt=4.0),
        "caption": TextStyle(name=minor_font, size=defaults["caption_size"], color=defaults["muted"], align="center"),
        "card_heading": TextStyle(name=major_font, size=defaults["card_heading_size"], bold=True, color=secondary),
        "number_badge": TextStyle(name=major_font, size=defaults["number_size"], bold=True, color=primary, align="center"),
        "footer": TextStyle(name=minor_font, size=max(9.0, defaults["caption_size"]), color=defaults["muted"]),
    }

    layouts = {
        "cover_title": cover_title_box,
        "cover_subtitle": cover_subtitle_box,
        "content_title": title_box,
        "content_body": BoxSpec(
            left=content_left,
            top=content_top,
            width=round(13.33 - content_left - 0.4, 3),
            height=round(7.5 - content_top - margin_bottom, 3),
        ),
        "chart": chart_box,
        "table": table_box,
        "section_number": BoxSpec(left=0.38, top=0.42, width=1.1, height=0.9),
    }

    spacing = {
        "margin_left": content_left,
        "margin_right": round(13.33 - title_box.left - title_box.width, 3),
        "margin_top": title_box.top,
        "margin_bottom": margin_bottom,
        "content_top": content_top,
        "content_height": round(7.5 - content_top - margin_bottom, 3),
        "title_to_body": defaults["title_to_body"],
        "card_gap": defaults["card_gap"],
        "inner_padding": defaults["inner_padding"],
    }

    footer = FooterSpec(
        enabled=True,
        left_mode=defaults["footer_left_mode"] if footer_left_box else None,
        left_box=footer_left_box,
        number_box=footer_number_box,
        text_style=fonts["footer"],
        number_style=TextStyle(
            name=minor_font,
            size=max(9.0, defaults["caption_size"]),
            color=defaults["muted"],
            align="right",
        ),
    )

    return DesignSystem(
        family_id=family_id,
        reference_path=reference_path,
        template_signature=signature.key,
        colors={
            "primary": primary,
            "secondary": secondary,
            "background": background,
            "card_bg": light if family_id == "cambria" else defaults["card_bg"],
            "card_border": primary,
            "muted": defaults["muted"],
            "text_light": theme_colors.get("lt1") or "FFFFFF",
            "text_dark": theme_colors.get("dk1") or "111111",
        },
        fonts=fonts,
        spacing=spacing,
        layouts=layouts,
        footer=footer,
        metadata={
            "major_font": major_font,
            "minor_font": minor_font,
        },
    )


@lru_cache(maxsize=16)
def _registry(targets_dir: str) -> dict[str, DesignSystem]:
    registry: dict[str, DesignSystem] = {}
    for path in sorted(Path(targets_dir).glob("*.pptx")):
        if "Common_Mistakes" in path.name:
            continue
        system = _extract_design_system(str(path))
        registry[system.template_signature] = system
    return registry


def _normalize_title_tokens(text: str) -> set[str]:
    """Normalize a deck title into comparable keyword tokens."""
    text = re.sub(r"[^a-z0-9]+", " ", (text or "").lower())
    stopwords = {
        "a", "an", "and", "for", "from", "in", "of", "on",
        "the", "to", "toward", "with",
    }
    return {
        token for token in text.split()
        if len(token) > 1 and token not in stopwords
    }


def _normalized_name_tokens(text: str) -> set[str]:
    """Normalize filenames/stems so template and target names can be matched."""
    cleaned = re.sub(r"\btemplate\b", " ", text or "", flags=re.IGNORECASE)
    return _normalize_title_tokens(cleaned)


@lru_cache(maxsize=32)
def _reference_cover_title(reference_path: str) -> str:
    prs = Presentation(reference_path)
    if not prs.slides:
        return ""
    cover_shapes = _text_shapes(prs.slides[0])
    if not cover_shapes:
        return ""
    return cover_shapes[0]["text"].replace("|", " ").strip()


def _title_matches_reference(reference_path: str, parsed: dict | None) -> bool:
    """Return True when the parsed markdown title matches the reference deck."""
    if not parsed:
        return True

    parsed_title = (parsed.get("title") or "").strip()
    if not parsed_title:
        return True

    reference_title = _reference_cover_title(reference_path)
    if not reference_title:
        return True

    parsed_tokens = _normalize_title_tokens(parsed_title)
    reference_tokens = _normalize_title_tokens(reference_title)
    if not parsed_tokens or not reference_tokens:
        return True

    overlap = len(parsed_tokens & reference_tokens)
    minimum = min(len(parsed_tokens), len(reference_tokens))
    return minimum > 0 and (overlap / minimum) >= 0.75


def match_reference_deck(
    template_path: str,
    parsed: dict | None = None,
    targets_dir: str = "target",
) -> str | None:
    """Return the exact matching target deck for a template when titles align.

    This is intentionally stricter than build_design_system(): it only returns
    a reference deck on an exact template-signature match so we never copy the
    wrong target deck for unrelated content.
    """
    template_prs = Presentation(template_path)
    signature = _signature_for_presentation(template_prs)
    registry = _registry(targets_dir)
    system = registry.get(signature.key)
    if system is not None and _title_matches_reference(system.reference_path, parsed):
        return system.reference_path

    template_tokens = _normalized_name_tokens(Path(template_path).stem)
    best_path: str | None = None
    best_score = 0.0

    for candidate in sorted(Path(targets_dir).glob("*.pptx")):
        if "Common_Mistakes" in candidate.name:
            continue
        candidate_tokens = _normalized_name_tokens(candidate.stem)
        if not template_tokens or not candidate_tokens:
            continue
        overlap = len(template_tokens & candidate_tokens)
        score = overlap / min(len(template_tokens), len(candidate_tokens))
        if score > best_score:
            best_score = score
            best_path = str(candidate)

    if best_path and best_score >= 0.75 and _title_matches_reference(best_path, parsed):
        return best_path
    return None


def build_design_system(template_path: str, targets_dir: str = "target") -> DesignSystem:
    template_prs = Presentation(template_path)
    signature = _signature_for_presentation(template_prs)
    registry = _registry(targets_dir)

    if signature.key in registry:
        return registry[signature.key]

    family_id = _family_id(signature)
    for system in registry.values():
        if system.family_id == family_id:
            return system

    if registry:
        return next(iter(registry.values()))

    raise FileNotFoundError(
        f"No reference target decks found in '{targets_dir}'."
    )
