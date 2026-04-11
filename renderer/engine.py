"""
renderer/engine.py — Main render loop: blueprint + template → output.pptx.

This is the Stage 3 entry point. The Renderer class:
  1. Selects the best template from the templates/ folder (or uses an explicit path)
  2. Loads the template as a Presentation object
  3. Records how many slides the template originally had
  4. Iterates over each slide in the blueprint and dispatches to the right renderer
  5. Removes the original template placeholder slides (keeping only our generated slides)
  6. Saves the final .pptx file to disk
"""

from __future__ import annotations

import logging
import os
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import config
from renderer.utils import (
    add_textbox,
    add_slide_title,
    add_slide_number,
    get_layout_by_name,
    get_blank_layout,
    remove_template_slides,
    strip_numeric_prefix,
    style_shape,
)
from renderer.layouts import render_content_slide   # handles all "content" slide types
from renderer.charts import render_chart_slide       # handles "chart" slides
from renderer.tables import render_table_slide       # handles "table" slides

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Public class
# ---------------------------------------------------------------------------

class Renderer:
    """Stage 3: Render a slide blueprint into a .pptx file.

    Instantiate once and call render() for each presentation you want to produce.

    Usage:
        renderer = Renderer(template_path=None, templates_dir="templates")
        renderer.render(blueprint, parsed, "outputs/result.pptx")
    """

    def __init__(
        self,
        template_path: str | None = None,
        templates_dir: str = "templates",
    ) -> None:
        """Store configuration; no I/O happens here.

        Args:
            template_path: Path to an explicit .pptx template file.
                           When None, the best template is auto-selected from
                           templates_dir based on content keyword matching.
            templates_dir: Folder containing available .pptx Slide Master files.
        """
        self.template_path = template_path
        self.templates_dir = templates_dir

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    def render(self, blueprint: dict, parsed: dict, output_path: str) -> None:
        """Render the slide blueprint to a .pptx file.

        This is the main method. It orchestrates the entire Stage 3 process:
        select template → load presentation → add slides → remove originals → save.

        Args:
            blueprint: Slide blueprint dict produced by Stage 2 (StorylineGenerator).
                       Contains a "slides" list with one dict per slide.
            parsed: Parsed document dict from Stage 1 (MarkdownParser).
                    Used only for template auto-selection heuristics.
            output_path: File path where the finished .pptx will be saved.
        """
        # Step 1: Decide which template to use
        template = self._select_template(parsed)
        logger.info("Using template: %s", template)

        # Step 2: Load the template — this copies all Slide Master layouts into memory
        prs = Presentation(template)

        # Step 2b: Extract the template's real theme colors and overwrite config
        # globals so every renderer module (layouts, charts, tables, visuals)
        # automatically uses the template's palette instead of the hardcoded red.
        self._apply_theme_colors(prs)

        # Step 3: Record the number of original template placeholder slides.
        # We'll use this count later to remove them after adding our content slides.
        original_slide_count = len(prs.slides)

        # Step 4: Iterate over each slide in the blueprint and render it
        slides_data = blueprint.get("slides", [])
        logger.info("Rendering %d slide(s) …", len(slides_data))

        for slide_data in slides_data:
            try:
                self._render_slide(prs, slide_data)
            except Exception as exc:
                # A single slide failing should never stop the whole presentation
                logger.warning(
                    "Slide %s (%s) failed: %s — skipping.",
                    slide_data.get("slide_number", "?"),
                    slide_data.get("type", "?"),
                    exc,
                )

        # Step 5: Delete the original template placeholder slides from the front
        # (all our new slides were appended after them)
        remove_template_slides(prs, original_slide_count)

        # Step 6: Save the finished presentation to disk
        os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
        prs.save(output_path)
        logger.info("Saved: %s", output_path)

    # ------------------------------------------------------------------
    # Template selection
    # ------------------------------------------------------------------

    def _select_template(self, parsed: dict) -> str:
        """Intelligently select the best template for the document.

        Strategy:
          1. If --template was passed, use it directly.
          2. Profile every .pptx in the templates folder (colors, mood, layouts).
          3. Use keyword heuristic (fast, deterministic, well-tuned).
          4. Only use LLM as tiebreaker when top-2 scores are within 10%.

        Args:
            parsed: Parsed document dict (title, sections, exec summary).

        Returns:
            Path string to the selected .pptx template.

        Raises:
            FileNotFoundError: If no templates exist.
        """
        if self.template_path:
            return self.template_path

        templates = sorted(Path(self.templates_dir).glob("*.pptx"))
        if not templates:
            raise FileNotFoundError(
                f"No .pptx templates found in '{self.templates_dir}'. "
                "Add at least one Slide Master .pptx file or use --template."
            )
        if len(templates) == 1:
            return str(templates[0])

        # Build a profile for each template (colors → mood)
        profiles = [self._profile_template(tpl) for tpl in templates]
        doc_summary = self._build_doc_summary(parsed)

        # Heuristic first (fast, deterministic, well-tuned keywords)
        idx, scores = self._heuristic_pick_template(profiles, doc_summary)

        # If top-2 scores are very close, use LLM as tiebreaker
        if scores is not None and len(scores) >= 2:
            sorted_scores = sorted(scores, reverse=True)
            top, second = sorted_scores[0], sorted_scores[1]
            if top > 0 and second / top > 0.9:
                try:
                    llm_idx = self._llm_pick_template(profiles, doc_summary)
                    if 0 <= llm_idx < len(templates):
                        logger.info(
                            "Auto-selected template: %s (LLM tiebreaker, heuristic was close)",
                            templates[llm_idx].name,
                        )
                        return str(templates[llm_idx])
                except Exception as exc:
                    logger.debug("LLM tiebreaker failed: %s — using heuristic winner.", exc)

        logger.info("Auto-selected template: %s (heuristic)", templates[idx].name)
        return str(templates[idx])

    # ------------------------------------------------------------------
    # Template profiling helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _profile_template(tpl_path: Path) -> dict:
        """Read a template and extract its visual identity.

        Returns a dict with filename, layout names, accent colors, and a
        human-readable mood string (e.g. "red/corporate, dark/formal").
        """
        prs = Presentation(str(tpl_path))
        layout_names = [l.name for l in prs.slide_layouts]

        accent_colors: dict[str, str] = {}
        theme_name = "Unknown"

        for rel in prs.part.rels.values():
            if "theme" not in str(rel.reltype).lower():
                continue
            theme_xml = rel.target_part.blob.decode("utf-8", errors="ignore")

            m = re.search(r'<a:theme[^>]*name="([^"]+)"', theme_xml)
            if m:
                theme_name = m.group(1)

            for tag, val in re.findall(
                r"<a:(\w+)>\s*<a:srgbClr val=\"([A-Fa-f0-9]{6})\"", theme_xml
            ):
                accent_colors[tag] = val
            break

        mood = Renderer._color_mood(accent_colors)

        return {
            "filename": tpl_path.name,
            "path": str(tpl_path),
            "layouts": layout_names,
            "theme_name": theme_name,
            "accent_colors": accent_colors,
            "mood": mood,
        }

    @staticmethod
    def _color_mood(accent_colors: dict[str, str]) -> str:
        """Derive a mood description from the primary accent colors only.

        Only accent1 and accent2 are considered — these define the template's
        dominant visual identity. Secondary accents (accent3–6) are ignored
        because most templates include a full rainbow that would blur the signal.
        """
        moods: list[str] = []
        for key in ("accent1", "accent2"):
            hex_val = accent_colors.get(key, "")
            if len(hex_val) != 6:
                continue
            r = int(hex_val[:2], 16)
            g = int(hex_val[2:4], 16)
            b = int(hex_val[4:6], 16)

            if r > 180 and g < 100 and b < 100:
                moods.append("red/corporate")
            elif r > 180 and g > 80 and b < 80:
                moods.append("orange/warm")
            elif g > max(r, b) and g > 80:
                moods.append("green/sustainability")
            elif b > max(r, g) and b > 120:
                moods.append("blue/technology")
            elif r < 60 and g < 60 and b < 60:
                moods.append("dark/formal")
        return ", ".join(dict.fromkeys(moods)) or "neutral/professional"

    @staticmethod
    def _build_doc_summary(parsed: dict) -> str:
        """Build a compact text summary for template matching."""
        title = parsed.get("title", "Untitled")
        subtitle = (parsed.get("subtitle", "") or "")[:120]
        headings = [s.get("heading", "") for s in parsed.get("sections", [])]
        exec_sum = (parsed.get("executive_summary", "") or "")[:250]
        return (
            f"Title: {title}\n"
            f"Subtitle: {subtitle}\n"
            f"Sections: {', '.join(headings)}\n"
            f"Summary: {exec_sum}"
        )

    def _llm_pick_template(
        self, profiles: list[dict], doc_summary: str
    ) -> int:
        """Use the LLM to pick the best template (one fast call).

        Returns a 0-based template index.
        Raises on any failure so the caller can fall back to heuristics.
        """
        from agents.base_agent import BaseAgent, MODEL  # noqa: PLC0415

        clients = BaseAgent._clients
        if not clients:
            raise ValueError("No Groq clients available for template selection.")

        idx = BaseAgent._client_index % len(clients)
        BaseAgent._client_index += 1
        client = clients[idx]

        # Build concise template descriptions for the LLM
        tpl_lines: list[str] = []
        for i, p in enumerate(profiles):
            accents = ", ".join(
                f"{k}=#{v}"
                for k, v in p["accent_colors"].items()
                if k.startswith("accent")
            )
            tpl_lines.append(
                f"Template {i + 1}: \"{p['filename']}\"  |  "
                f"Mood: {p['mood']}  |  "
                f"Layouts: {', '.join(p['layouts'])}  |  "
                f"Colors: {accents}"
            )

        prompt = (
            "Pick the single best PowerPoint template for this presentation.\n\n"
            f"DOCUMENT:\n{doc_summary}\n\n"
            f"TEMPLATES:\n" + "\n".join(tpl_lines) + "\n\n"
            "Match the document's topic to the template whose color mood fits best.\n"
            f"Reply with ONLY a number from 1 to {len(profiles)}. No explanation."
        )

        response = client.chat.completions.create(
            model=MODEL,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are a design expert. Pick the best matching "
                        "template. Reply with only the number."
                    ),
                },
                {"role": "user", "content": prompt},
            ],
            temperature=0.1,
            max_tokens=5,
        )
        text = (response.choices[0].message.content or "").strip()
        match = re.search(r"\d+", text)
        if match:
            return int(match.group()) - 1  # 1-based → 0-based
        raise ValueError(f"Could not parse template number from LLM: {text}")

    @staticmethod
    def _heuristic_pick_template(
        profiles: list[dict], doc_summary: str
    ) -> tuple[int, list[int]]:
        """Score templates against document keywords (no LLM needed).

        Returns:
            Tuple of (best_index, list_of_scores_per_template).
        """
        doc_lower = doc_summary.lower()

        # Strong-signal keywords that UNIQUELY identify a topic.
        # Ambiguous words (investment, market, strategy, growth, tech) are
        # excluded because they appear across domains (e.g. "tech acquisitions"
        # is corporate, not tech).
        GREEN_KW = {
            "sustainability", "environment", "renewable", "solar",
            "wind", "climate", "carbon", "emission", "biodiversity",
            "esg", "conservation", "photovoltaic", "decarbonization",
        }
        TECH_KW = {
            "ai", "technology", "digital", "software", "cyber", "cloud",
            "automation", "innovation", "machine", "blockchain", "saas",
            "infrastructure", "algorithm", "neural", "computing",
            "artificial", "bubble",
        }
        CORP_KW = {
            "acquisition", "merger", "consulting", "enterprise", "portfolio",
            "corporate", "management", "revenue", "valuation",
            "stakeholder", "governance", "compliance", "fiscal",
            "realignment", "reinvention",
        }

        # Use whole-word matching (prevents "ai" matching inside "sustainability")
        doc_words = set(re.findall(r"\b[a-z0-9]+\b", doc_lower))

        # Split title from body — title matches are a stronger signal
        title_line = doc_lower.split("\n")[0] if "\n" in doc_lower else doc_lower
        title_words = set(re.findall(r"\b[a-z0-9]+\b", title_line))

        # Count keyword OCCURRENCES in full text (not just presence) for body,
        # so a doc that mentions "acquisition" 30 times outscores one that
        # mentions it once.
        def _freq_score(keywords: set[str], text: str) -> int:
            return sum(
                len(re.findall(r"\b" + re.escape(kw) + r"\b", text))
                for kw in keywords
            )

        best_idx = 0
        best_score = -1
        all_scores: list[int] = []

        for i, p in enumerate(profiles):
            score = 0
            mood = p.get("mood", "").lower()

            if "green" in mood or "sustainability" in mood:
                kw = GREEN_KW
            elif "blue" in mood or "technology" in mood:
                kw = TECH_KW
            elif "orange" in mood or "warm" in mood:
                kw = CORP_KW
            else:
                kw = CORP_KW | TECH_KW

            # Title matches (strong signal)
            score += len(kw & title_words) * 20
            # Body frequency (volume signal)
            score += _freq_score(kw, doc_lower)

            all_scores.append(score)
            if score > best_score:
                best_score = score
                best_idx = i

        return best_idx, all_scores

    # ------------------------------------------------------------------
    # Theme color application
    # ------------------------------------------------------------------

    @staticmethod
    def _apply_theme_colors(prs: Presentation) -> None:
        """Extract the template's accent colors and overwrite config globals.

        This makes every reference to config.COLOR_PRIMARY, COLOR_CARD_BG,
        COLOR_CARD_BORDER, COLOR_HEADER_BG, and CHART_COLORS automatically
        use the template's real palette instead of the hardcoded red defaults.
        """
        from lxml import etree  # noqa: PLC0415

        try:
            master = prs.slide_masters[0]
            theme_part = master.part.part_related_by(
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
            )
            root = etree.fromstring(theme_part.blob)
            ns = "http://schemas.openxmlformats.org/drawingml/2006/main"

            color_map: dict[str, RGBColor] = {}
            for elem in root.iter(f"{{{ns}}}srgbClr"):
                val = elem.get("val", "")
                if len(val) == 6:
                    try:
                        r, g, b = int(val[:2], 16), int(val[2:4], 16), int(val[4:6], 16)
                        parent_tag = elem.getparent().tag.split("}")[-1]
                        color_map[parent_tag] = RGBColor(r, g, b)
                    except ValueError:
                        pass

            if not color_map:
                return

            raw_accent1 = color_map.get("accent1")
            accent2 = color_map.get("accent2")

            # If accent1 is very light (a tint/background color), promote accent2
            # to primary — e.g. template 3 has accent1=#EFF3E5 (pastel) but
            # accent2=#33621A (the actual brand green).
            if raw_accent1 and (raw_accent1[0] + raw_accent1[1] + raw_accent1[2]) / 3 > 200:
                primary = accent2 or raw_accent1
            else:
                primary = raw_accent1 or config.COLOR_PRIMARY

            if accent2 is None:
                accent2 = primary
            dk1 = color_map.get("dk1", config.COLOR_TEXT_DARK)
            lt1 = color_map.get("lt1", config.COLOR_TEXT_LIGHT)

            # Derive card background: a very light tint of the primary color
            pr, pg, pb = primary[0], primary[1], primary[2]
            card_bg = RGBColor(
                min(255, 230 + pr // 10),
                min(255, 230 + pg // 10),
                min(255, 230 + pb // 10),
            )

            # Overwrite the config module globals in-place
            config.COLOR_PRIMARY = primary
            config.COLOR_CARD_BORDER = primary
            config.COLOR_HEADER_BG = primary
            config.COLOR_CARD_BG = card_bg
            config.COLOR_TEXT_DARK = dk1
            config.COLOR_TEXT_LIGHT = lt1

            # Rebuild chart colors from the template palette
            config.CHART_COLORS = [
                primary,
                dk1,
                card_bg,
                RGBColor(0x66, 0x66, 0x66),
                accent2,
                RGBColor(
                    max(0, pr - 40), max(0, pg - 40), max(0, pb - 40),
                ),
            ]

            logger.info(
                "Theme colors applied: primary=#%02X%02X%02X, accent2=#%02X%02X%02X",
                primary[0], primary[1], primary[2],
                accent2[0], accent2[1], accent2[2],
            )

        except Exception as exc:
            logger.debug("Could not extract theme colors: %s — using defaults.", exc)

    # ------------------------------------------------------------------
    # Slide dispatch
    # ------------------------------------------------------------------

    def _render_slide(self, prs: Presentation, slide_data: dict) -> None:
        """Route a single slide to the correct renderer based on its type.

        Special slide types (cover, section_divider, thank_you) have their own
        methods because they use specific template layouts and different designs.
        All other types (content, chart, table, agenda, executive_summary,
        conclusion) go through the generic content pipeline.

        Args:
            prs: The Presentation object being built.
            slide_data: Single slide dict from the blueprint JSON.
        """
        slide_type = slide_data.get("type", "content")
        slide_num = slide_data.get("slide_number", 0)

        if slide_type == "cover":
            # Uses the template's dedicated Cover layout (has background image/branding)
            self._render_cover(prs, slide_data)

        elif slide_type == "section_divider":
            # Uses the template's Divider layout for section break slides
            self._render_divider(prs, slide_data)

        elif slide_type == "thank_you":
            # Uses the template's Thank You layout for the closing slide
            self._render_thank_you(prs, slide_data)

        elif slide_type == "chart":
            # Chart slides use a Blank layout so we have full control over positioning
            layout = get_blank_layout(prs)
            slide = prs.slides.add_slide(layout)
            self._clear_placeholders(slide)   # remove any placeholder shapes from the layout
            render_chart_slide(slide, slide_data, slide_num)

        elif slide_type == "table":
            # Table slides also use Blank layout
            layout = get_blank_layout(prs)
            slide = prs.slides.add_slide(layout)
            self._clear_placeholders(slide)
            render_table_slide(slide, slide_data, slide_num)

        else:
            # All remaining types (content, agenda, executive_summary, conclusion)
            # use the Blank layout and the layouts.py dispatcher
            layout = get_blank_layout(prs)
            slide = prs.slides.add_slide(layout)
            self._clear_placeholders(slide)
            render_content_slide(slide, slide_data, slide_num)

    # ------------------------------------------------------------------
    # Specific slide type renderers
    # ------------------------------------------------------------------

    def _render_cover(self, prs: Presentation, data: dict) -> None:
        """Render the cover/title slide.

        Tries the template's Cover layout first. If it has usable placeholders,
        fills them with the title and subtitle. Otherwise falls back to
        manually positioned text boxes drawn on top of the background.
        """
        # Look for a layout named "cover" — common in most Slide Masters
        layout = get_layout_by_name(prs, "cover")
        # Some templates name the cover layout with a "1_" prefix — try that too
        if "cover" not in layout.name.lower():
            layout = get_layout_by_name(prs, "1_")

        slide = prs.slides.add_slide(layout)

        title = data.get("title", "")
        subtitle = data.get("subtitle", "")

        # Hard cap: cover subtitle must be max 20 words — never dump exec summary here
        words = subtitle.split()
        if len(words) > 20:
            subtitle = " ".join(words[:20]) + "…"

        # Attempt to fill the template's built-in title/subtitle placeholders
        filled = self._fill_placeholders(slide, title=title, subtitle=subtitle)

        if not filled:
            # Placeholders didn't work — draw text boxes manually
            self._clear_placeholders(slide)
            # Large centred title in white (template backgrounds are usually dark)
            add_textbox(
                slide, title,
                left=config.MARGIN_LEFT,
                top=Inches(2.5),           # roughly vertically centred on a 7.5" slide
                width=config.CONTENT_WIDTH,
                height=Inches(1.5),
                font_size=Pt(40), bold=True,
                color=config.COLOR_TEXT_LIGHT,
                align=PP_ALIGN.CENTER,
                font_name=config.TITLE_FONT,
            )
            if subtitle:
                add_textbox(
                    slide, subtitle,
                    left=config.MARGIN_LEFT,
                    top=Inches(4.2),       # below the title
                    width=config.CONTENT_WIDTH,
                    height=Inches(0.8),
                    font_size=config.SUBTITLE_FONT_SIZE,
                    color=config.COLOR_TEXT_LIGHT,
                    align=PP_ALIGN.CENTER,
                )

    def _render_divider(self, prs: Presentation, data: dict) -> None:
        """Render a section divider slide.

        Full-bleed red background with a large serif section number on the
        left, the section title below it, and a subtle subtitle underneath.
        This replaces the old dark-bg + dash approach to match the target
        deck's rhythm of "every new section gets a red page".
        """
        # Always use Blank layout — we draw everything ourselves
        layout = get_blank_layout(prs)
        slide = prs.slides.add_slide(layout)
        self._clear_placeholders(slide)

        title = strip_numeric_prefix(data.get("title", ""))
        subtitle = data.get("subtitle", "")
        section_num = data.get("section_number", "")

        # Full-bleed red background
        bg = slide.shapes.add_shape(
            1, 0, 0, config.SLIDE_WIDTH, config.SLIDE_HEIGHT,
        )
        style_shape(bg, fill_color=config.COLOR_PRIMARY, line_color=None)

        # Large section number — huge serif outline on the left
        if section_num:
            add_textbox(
                slide, section_num,
                left=config.MARGIN_LEFT, top=Inches(1.4),
                width=Inches(3.0), height=Inches(1.6),
                font_size=Pt(80), bold=True,
                color=config.COLOR_TEXT_LIGHT,
                font_name=config.TITLE_FONT,
            )

        # Thin white divider line under the number
        line_top = Inches(3.1) if section_num else Inches(2.2)
        line = slide.shapes.add_shape(
            1, config.MARGIN_LEFT, line_top,
            Inches(1.2), Inches(0.035),
        )
        style_shape(line, fill_color=config.COLOR_TEXT_LIGHT, line_color=None)

        # Section title — large, white, serif
        title_top = line_top + Inches(0.25)
        add_textbox(
            slide, title,
            left=config.MARGIN_LEFT, top=title_top,
            width=config.CONTENT_WIDTH, height=Inches(1.5),
            font_size=Pt(36), bold=True,
            color=config.COLOR_TEXT_LIGHT,
            font_name=config.TITLE_FONT,
        )

        # Subtitle — light pink for readability on red
        if subtitle:
            from pptx.dml.color import RGBColor
            add_textbox(
                slide, subtitle,
                left=config.MARGIN_LEFT, top=title_top + Inches(1.6),
                width=config.CONTENT_WIDTH, height=Inches(0.8),
                font_size=config.SUBTITLE_FONT_SIZE,
                color=RGBColor(0xFF, 0xE5, 0xE1),
            )

    def _render_thank_you(self, prs: Presentation, data: dict) -> None:
        """Render the closing Thank You slide.

        Tries the template's Thank You layout first, then falls back to
        a manually drawn centred layout on a default background.
        """
        # Look for a layout whose name contains "thank"
        layout = get_layout_by_name(prs, "thank")
        slide = prs.slides.add_slide(layout)

        title = data.get("title", "Thank You")
        subtitle = data.get("subtitle", "")

        filled = self._fill_placeholders(slide, title=title, subtitle=subtitle)

        if not filled:
            self._clear_placeholders(slide)
            # Large centred "Thank You" text
            add_textbox(
                slide, title,
                left=config.MARGIN_LEFT, top=Inches(2.8),
                width=config.CONTENT_WIDTH, height=Inches(1.2),
                font_size=Pt(48), bold=True,
                color=config.COLOR_TEXT_LIGHT,
                align=PP_ALIGN.CENTER,
                font_name=config.TITLE_FONT,
            )
            if subtitle:
                add_textbox(
                    slide, subtitle,
                    left=config.MARGIN_LEFT, top=Inches(4.2),
                    width=config.CONTENT_WIDTH, height=Inches(0.6),
                    font_size=config.SUBTITLE_FONT_SIZE,
                    color=config.COLOR_TEXT_MUTED,
                    align=PP_ALIGN.CENTER,
                )

    # ------------------------------------------------------------------
    # Placeholder helpers
    # ------------------------------------------------------------------

    def _fill_placeholders(
        self, slide, title: str = "", subtitle: str = ""
    ) -> bool:
        """Try to populate the slide's built-in layout placeholders.

        Template layouts often have pre-positioned title (idx=0) and body/subtitle
        (idx=1) placeholders that inherit the template's font and position styling.
        Using them is preferable to drawing new text boxes from scratch.

        Args:
            slide: python-pptx slide object.
            title: Text for the title placeholder (idx=0).
            subtitle: Text for the subtitle/body placeholder (idx=1).

        Returns:
            True if at least the title placeholder was successfully populated.
            False means the caller should fall back to manual text boxes.
        """
        filled = False
        for ph in slide.placeholders:
            try:
                idx = ph.placeholder_format.idx
                if idx == 0 and title:
                    ph.text = title      # title placeholder
                    filled = True
                elif idx == 1 and subtitle:
                    ph.text = subtitle   # subtitle/body placeholder
            except Exception as exc:
                logger.debug("Placeholder fill failed (idx=%s): %s", idx, exc)
        return filled

    def _clear_placeholders(self, slide) -> None:
        """Remove all placeholder shapes from a slide.

        When we draw shapes manually (not using template placeholders), any
        existing placeholder shapes from the layout would sit underneath our
        custom shapes and could show default "Click to add title" prompt text.
        Removing them gives us a truly blank canvas.

        Args:
            slide: python-pptx slide object.
        """
        # Build a list first — modifying the collection while iterating causes issues
        placeholders = list(slide.placeholders)
        for ph in placeholders:
            try:
                sp = ph._element             # the underlying <p:sp> XML element
                sp.getparent().remove(sp)    # detach it from the slide's shape tree
            except Exception:
                pass   # if removal fails for any placeholder, just skip it
