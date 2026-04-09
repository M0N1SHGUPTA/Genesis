"""
renderer/infographics.py — Shape-based infographic renderers.

All visuals are built exclusively from native python-pptx shapes
(rectangles, ovals, text boxes). No external images, SVGs, or icons.

These functions are called by layouts.py when a slide needs a more complex
visual that goes beyond what simple text boxes provide:
  - render_vertical_timeline  — vertical variant of the timeline
  - render_wrapped_process_flow — two-row process flow for 6+ steps
  - render_comparison_grid — 2×2 grid of comparison boxes
"""

from __future__ import annotations

import logging

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import config
from renderer.utils import add_textbox, style_shape

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Vertical timeline
# Used when the number of steps exceeds what fits on a horizontal layout,
# or when the calling layout explicitly requests a vertical orientation.
# ---------------------------------------------------------------------------

def render_vertical_timeline(slide, steps: list[dict], left: float, top: float,
                               width: float, height: float) -> None:
    """Render a vertical timeline: circles on the left connected by a vertical line.

    Each step has a circle (with number) on the left side and heading +
    description text to the right.

    Args:
        slide: python-pptx slide object.
        steps: List of dicts, each with keys: number, heading, description.
        left, top, width, height: Bounding box for the entire infographic (EMU).
    """
    if not steps:
        return   # nothing to draw

    n = len(steps)
    row_height = height / n              # divide available height equally among steps
    circle_size = Inches(0.45)           # diameter of each numbered circle
    line_x = left + circle_size / 2      # x-coordinate of the vertical connector line (centred on circles)

    # Draw the vertical connector line that links all circles
    # It spans the full height and is centred on the circles' x-position
    line = slide.shapes.add_shape(
        1,   # rectangle shape type
        int(line_x - Inches(0.02)), int(top),   # slight left offset to visually centre the line
        int(Inches(0.04)), int(height),          # 0.04" wide (thin bar acting as a line)
    )
    style_shape(line, fill_color=config.COLOR_PRIMARY, line_color=None)

    # Text content starts to the right of the circles with a small gap
    text_left = left + circle_size + Inches(0.2)
    text_width = width - circle_size - Inches(0.2)

    for i, step in enumerate(steps):
        row_top = top + i * row_height
        # Vertically centre the circle within its allocated row
        circle_top = row_top + (row_height - circle_size) / 2

        # Red filled circle sitting on the vertical line
        circle = slide.shapes.add_shape(
            9,   # shape type 9 = oval (creates a circle when width == height)
            int(left), int(circle_top),
            int(circle_size), int(circle_size),
        )
        style_shape(circle, fill_color=config.COLOR_PRIMARY, line_color=None)

        # White step number inside the circle
        number = step.get("number", str(i + 1).zfill(2))
        add_textbox(
            slide, number,
            left=int(left), top=int(circle_top),
            width=int(circle_size), height=int(circle_size),
            font_size=Pt(12), bold=True, color=config.COLOR_TEXT_LIGHT,
            align=PP_ALIGN.CENTER,
        )

        # Bold heading to the right of the circle
        heading = step.get("heading", "")
        add_textbox(
            slide, heading,
            left=int(text_left), top=int(row_top + Inches(0.05)),
            width=int(text_width), height=int(Inches(0.35)),
            font_size=Pt(13), bold=True, color=config.COLOR_TEXT_DARK,
        )

        # Muted description text below the heading
        desc = step.get("description", "")
        add_textbox(
            slide, desc,
            left=int(text_left), top=int(row_top + Inches(0.42)),
            width=int(text_width), height=int(row_height - Inches(0.5)),
            font_size=Pt(11), color=config.COLOR_TEXT_MUTED,
        )


# ---------------------------------------------------------------------------
# Wrapped process flow (two rows)
# Used when a process has more steps than fit on a single horizontal row (>5).
# Splits steps across two rows: first 5 on row 1, remainder on row 2.
# ---------------------------------------------------------------------------

def render_wrapped_process_flow(slide, steps: list[dict], left: float, top: float,
                                  width: float, height: float) -> None:
    """Render a two-row process flow for workflows with more than 5 steps.

    Each row contains boxes connected by arrow bars. The second row
    continues the sequence from where the first row left off.

    Args:
        slide: python-pptx slide object.
        steps: List of dicts, each with keys: number, heading, description.
        left, top, width, height: Bounding box for the entire infographic (EMU).
    """
    if not steps:
        return

    # Split steps: first row gets up to 5, second row gets the rest
    row1 = steps[:5]
    row2 = steps[5:]
    rows = [row1] + ([row2] if row2 else [])   # only include row2 if it has content

    # Divide available height between the rows, leaving a gap between them
    row_height = height / len(rows) - config.ELEMENT_GAP
    arrow_w = Inches(0.25)   # width of each arrow bar between boxes

    for r_idx, row_steps in enumerate(rows):
        # Each row starts at top + (row index × row height + gap)
        row_top = top + r_idx * (row_height + config.ELEMENT_GAP)
        n = len(row_steps)
        total_arrow = (n - 1) * arrow_w                 # total space used by all arrows in this row
        box_width = (width - total_arrow) / n            # equal width per box

        for i, step in enumerate(row_steps):
            box_left = left + i * (box_width + arrow_w)

            # Cream card box with red border
            box = slide.shapes.add_shape(
                1, int(box_left), int(row_top), int(box_width), int(row_height)
            )
            style_shape(box, fill_color=config.COLOR_CARD_BG, line_color=config.COLOR_PRIMARY)

            # Large red step number at the top of the box
            number = step.get("number", str(i + 1).zfill(2))
            add_textbox(
                slide, number,
                left=int(box_left + config.INNER_PADDING),
                top=int(row_top + config.INNER_PADDING),
                width=int(box_width - 2 * config.INNER_PADDING),
                height=int(Inches(0.35)),
                font_size=Pt(18), bold=True, color=config.COLOR_PRIMARY,
                align=PP_ALIGN.CENTER,
            )

            # Step heading below the number
            heading = step.get("heading", "")
            add_textbox(
                slide, heading,
                left=int(box_left + config.INNER_PADDING),
                top=int(row_top + Inches(0.45)),
                width=int(box_width - 2 * config.INNER_PADDING),
                height=int(Inches(0.4)),
                font_size=Pt(12), bold=True, color=config.COLOR_TEXT_DARK,
                align=PP_ALIGN.CENTER,
            )

            # Draw a thin red arrow bar between this box and the next
            # (skip after the last box in each row)
            if i < n - 1:
                arrow_left = box_left + box_width
                arrow = slide.shapes.add_shape(
                    1,
                    int(arrow_left),
                    int(row_top + row_height / 2 - Inches(0.03)),   # vertically centred
                    int(arrow_w), int(Inches(0.06)),
                )
                style_shape(arrow, fill_color=config.COLOR_PRIMARY, line_color=None)


# ---------------------------------------------------------------------------
# Comparison grid (up to 4 items in a 2×2 layout)
# Each cell has a contrasting background color, heading, and bullet points.
# ---------------------------------------------------------------------------

def render_comparison_grid(slide, items: list[dict], left: float, top: float,
                             width: float, height: float) -> None:
    """Render a 2×2 grid of comparison boxes (up to 4 items).

    Items are laid out in a grid: 1-2 items in a single row, 3-4 items
    in two rows. Each box uses a different accent color.

    Args:
        slide: python-pptx slide object.
        items: List of dicts, each with keys: heading, points.
        left, top, width, height: Bounding box for the grid (EMU).
    """
    items = items[:4]   # maximum 4 items fit in a 2×2 grid
    if not items:
        return

    n = len(items)
    # Use 2 columns for 3-4 items, 1 column for 1-2 items
    cols = 2 if n > 2 else n
    rows = (n + cols - 1) // cols   # ceiling division to get row count

    # Calculate individual cell dimensions
    cell_width = (width - (cols - 1) * config.CARD_GAP) / cols
    cell_height = (height - (rows - 1) * config.CARD_GAP) / rows

    # Rotate through these background colors for visual variety
    colors = [
        config.COLOR_PRIMARY,      # red (item 0)
        config.COLOR_TEXT_DARK,    # near-black (item 1)
        config.COLOR_CARD_BG,      # light cream (item 2)
        config.COLOR_TEXT_MUTED,   # gray (item 3)
    ]

    for i, item in enumerate(items):
        # Calculate grid position for this item
        row = i // cols
        col = i % cols
        cell_left = left + col * (cell_width + config.CARD_GAP)
        cell_top = top + row * (cell_height + config.CARD_GAP)

        bg_color = colors[i % len(colors)]

        # Light backgrounds need dark text; dark backgrounds need light text
        text_color = (
            config.COLOR_TEXT_LIGHT
            if bg_color not in (config.COLOR_CARD_BG, config.COLOR_TEXT_MUTED)
            else config.COLOR_TEXT_DARK
        )

        # Solid colored background box (no border)
        box = slide.shapes.add_shape(
            1, int(cell_left), int(cell_top), int(cell_width), int(cell_height)
        )
        style_shape(box, fill_color=bg_color, line_color=None)

        # Heading at the top of the cell
        heading = item.get("heading", "")
        if heading:
            add_textbox(
                slide, heading,
                left=int(cell_left + config.INNER_PADDING),
                top=int(cell_top + config.INNER_PADDING),
                width=int(cell_width - 2 * config.INNER_PADDING),
                height=int(Inches(0.4)),
                font_size=config.CARD_HEADING_SIZE, bold=True,
                color=text_color, align=PP_ALIGN.CENTER,
            )

        # Bullet points below the heading
        points = item.get("points", [])
        if points:
            from renderer.utils import add_bullet_textbox
            add_bullet_textbox(
                slide, points,
                left=int(cell_left + config.INNER_PADDING),
                top=int(cell_top + Inches(0.55)),
                width=int(cell_width - 2 * config.INNER_PADDING),
                height=int(cell_height - Inches(0.65)),
                font_size=Pt(12), color=text_color,
            )
