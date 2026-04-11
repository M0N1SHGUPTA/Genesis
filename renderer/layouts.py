"""
renderer/layouts.py — Layout rendering functions for content slides.

Each function receives a slide object + slide_data dict and draws all
shapes at precise grid-aligned positions using constants from config.py.

How it works:
  render_content_slide() reads the "layout" field from the blueprint and
  dispatches to the correct private function (_two_column, _three_cards, etc.).
  Agenda slides get their own dedicated layout.
  If a layout fails, it silently falls back to _single_focus.

Layouts implemented:
  two_column, three_cards, key_stats, timeline,
  process_flow, comparison, icon_list, single_focus, agenda
"""

from __future__ import annotations

import logging

from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import config
from renderer.utils import (
    add_textbox,
    add_bullet_textbox,
    add_slide_number,
    add_slide_title,
    pick_contrasting_text,
    style_shape,
)
from renderer.visuals import (
    draw_card_with_divider,
    draw_icon_glyph,
    draw_numbered_badge,
    draw_red_full_background,
    draw_red_left_sidebar,
    draw_red_top_pill,
    icon_for_text,
)

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Short aliases for frequently used config values
# ---------------------------------------------------------------------------
ML = config.MARGIN_LEFT     # left edge of the usable content area
CT = config.CONTENT_TOP     # top edge of the body area (below title)
CW = config.CONTENT_WIDTH   # total usable width
CH = config.CONTENT_HEIGHT  # total usable height (below title, above bottom margin)
GAP = config.ELEMENT_GAP    # vertical gap between stacked elements
PAD = config.INNER_PADDING  # padding inside cards/boxes
CGAP = config.CARD_GAP      # horizontal gap between side-by-side cards


def _truncate(text: str, max_words: int) -> str:
    """Trim text to max_words words, splitting on clause boundaries when possible.

    Instead of blindly chopping at word N, this tries to find a natural break
    (comma, semicolon, colon, em-dash) within the first max_words words.
    Falls back to word-count truncation if no clause boundary exists.
    """
    words = text.split()
    if len(words) <= max_words:
        return text
    # Try to find a clause boundary within the first max_words words
    fragment = " ".join(words[:max_words])
    for sep in (",", ";", ":", " –", " —", " -"):
        idx = fragment.rfind(sep)
        # Only accept if the split keeps at least 3 words of substance
        if idx > 0 and len(fragment[:idx].split()) >= 3:
            return fragment[:idx].rstrip(".,;: –—-")
    return fragment + "…"


def _heading_truncate(text: str, max_words: int = 6) -> str:
    """Truncate text for a card/step heading and title-case the result.

    Headings should be short, punchy, and title-cased. This runs
    _truncate then applies title-casing for a polished look.
    """
    result = _truncate(text, max_words)
    # Don't title-case if it's already an acronym-heavy string (e.g. "UAE AI GDP")
    upper_count = sum(1 for w in result.split() if w.isupper() and len(w) > 1)
    if upper_count > len(result.split()) / 2:
        return result
    return result.title()


# ---------------------------------------------------------------------------
# Public dispatcher
# ---------------------------------------------------------------------------

def render_content_slide(slide, slide_data: dict, slide_num: int) -> None:
    """Entry point for all non-special content slides.

    Reads the "layout" key from slide_data, adds the title,
    then calls the matching private layout function.

    Args:
        slide: python-pptx slide object (already added to the presentation).
        slide_data: Single slide dict from the blueprint JSON.
        slide_num: Slide number shown in the bottom-right corner.
    """
    layout = slide_data.get("layout", "single_focus")
    title = slide_data.get("title", "")
    slide_type = slide_data.get("type", "content")

    # Layouts that paint their own title/sidebar and must NOT have the
    # standard top title bar layered on top of them.
    _FULL_BLEED_LAYOUTS = {"exec_summary_with_photo", "two_col_sidebar", "five_cards_row"}
    needs_standard_title = layout not in _FULL_BLEED_LAYOUTS

    if needs_standard_title:
        # Title is added first so all layout functions can assume it's already there
        add_slide_title(slide, title)

    # Agenda gets its own dedicated renderer regardless of layout field
    if slide_type == "agenda":
        try:
            _agenda_layout(slide, slide_data)
        except Exception as exc:
            logger.warning("Agenda layout error: %s", exc)
            _single_focus(slide, slide_data)
        add_slide_number(slide, slide_num)
        return

    # Map layout name strings to their render functions
    dispatch = {
        "two_column":               _two_column,
        "three_cards":              _three_cards,
        "key_stats":                _key_stats,
        "timeline":                 _timeline,
        "process_flow":             _process_flow,
        "comparison":               _comparison,
        "icon_list":                _icon_list,
        "single_focus":             _single_focus,
        "six_cards":                _six_cards,
        "five_cards_row":           _five_cards_row,
        "two_col_sidebar":          _two_col_sidebar,
        "exec_summary_with_photo":  _exec_summary_with_photo,
    }

    fn = dispatch.get(layout, _single_focus)
    try:
        fn(slide, slide_data)
    except Exception as exc:
        logger.warning("Layout '%s' render error: %s — falling back to single_focus", layout, exc)
        try:
            _single_focus(slide, slide_data)
        except Exception:
            pass

    add_slide_number(slide, slide_num)


# ---------------------------------------------------------------------------
# Layout: agenda
# Two-column list of all section headings, numbered.
# ---------------------------------------------------------------------------

def _agenda_layout(slide, data: dict) -> None:
    """Render an agenda/TOC slide.

    Splits items into two columns when there are more than 4 items.
    Uses numbered entries instead of bullet dots.

    Blueprint key expected:
        points  (list[str]) — list of section headings
    """
    points = data.get("points", [])
    if not points:
        return

    # Cap to avoid overflow — 12 items max (6 per column)
    points = points[:12]

    ITEM_FONT = Pt(15)
    NUM_FONT = Pt(18)
    ITEM_HEIGHT = Inches(0.65)
    NUM_WIDTH = Inches(0.55)
    badge_text_color = pick_contrasting_text(config.COLOR_PRIMARY)

    if len(points) <= 4:
        # Single column, large numbered items
        for i, point in enumerate(points):
            row_top = CT + i * (ITEM_HEIGHT + Inches(0.1))
            # Red number badge
            num_box = slide.shapes.add_shape(9, ML, row_top + Inches(0.05),
                                             Inches(0.45), Inches(0.45))
            style_shape(num_box, fill_color=config.COLOR_PRIMARY, line_color=None)
            add_textbox(
                slide, str(i + 1).zfill(2),
                left=ML, top=row_top + Inches(0.05),
                width=Inches(0.45), height=Inches(0.45),
                font_size=NUM_FONT, bold=True,
                color=badge_text_color,
                align=PP_ALIGN.CENTER,
            )
            # Item text
            add_textbox(
                slide, _truncate(point, 12),
                left=ML + Inches(0.6), top=row_top,
                width=CW - Inches(0.6), height=ITEM_HEIGHT,
                font_size=ITEM_FONT, color=config.COLOR_TEXT_DARK,
            )
    else:
        # Two columns — split items in half
        mid = (len(points) + 1) // 2
        left_pts = points[:mid]
        right_pts = points[mid:]

        col_width = (CW - CGAP) / 2
        col_height = CH - GAP

        for col_i, pts in enumerate([left_pts, right_pts]):
            col_left = ML + col_i * (col_width + CGAP)

            # Subtle card background
            card = slide.shapes.add_shape(1, col_left, CT, col_width, col_height)
            style_shape(card, fill_color=config.COLOR_CARD_BG, line_color=config.COLOR_CARD_BORDER)

            for row_i, point in enumerate(pts[:6]):
                global_num = col_i * mid + row_i + 1
                row_top = CT + PAD + row_i * (ITEM_HEIGHT + Inches(0.05))

                # Red number
                add_textbox(
                    slide, str(global_num).zfill(2),
                    left=col_left + PAD, top=row_top,
                    width=NUM_WIDTH, height=ITEM_HEIGHT,
                    font_size=Pt(14), bold=True,
                    color=config.COLOR_PRIMARY,
                    align=PP_ALIGN.CENTER,
                )
                # Heading text
                add_textbox(
                    slide, _truncate(point, 10),
                    left=col_left + PAD + NUM_WIDTH + Inches(0.1),
                    top=row_top,
                    width=col_width - PAD - NUM_WIDTH - Inches(0.2),
                    height=ITEM_HEIGHT,
                    font_size=Pt(13),
                    color=config.COLOR_TEXT_DARK,
                )


# ---------------------------------------------------------------------------
# Layout: two_column
# Splits the content area into two equal side-by-side cards.
# ---------------------------------------------------------------------------

def _two_column(slide, data: dict) -> None:
    """Render a two-column layout.

    Blueprint keys expected:
        left.heading  (str)
        left.points   (list[str])
        right.heading (str)
        right.points  (list[str])
    """
    col_width = (CW - CGAP) / 2
    col_height = CH - GAP

    for i, side in enumerate(("left", "right")):
        col = data.get(side, {})
        if not isinstance(col, dict):
            col = {}

        left = ML + i * (col_width + CGAP)

        # Cream card background with red border
        card = slide.shapes.add_shape(1, left, CT, col_width, col_height)
        style_shape(card, fill_color=config.COLOR_CARD_BG, line_color=config.COLOR_CARD_BORDER)

        # Red accent top bar (visual heading indicator)
        accent = slide.shapes.add_shape(1, left, CT, col_width, Inches(0.07))
        style_shape(accent, fill_color=config.COLOR_PRIMARY, line_color=None)

        # Icon glyph next to the heading — derived from heading text.
        # Guide rule: "Icons are rarely used" → add icons to all category layouts.
        heading = col.get("heading", "")
        icon_name = col.get("icon") or icon_for_text(heading)
        icon_size = Inches(0.35)
        draw_icon_glyph(
            slide, icon_name,
            left=left + PAD, top=CT + Inches(0.16),
            size=icon_size,
        )

        # Heading inside card (shifted right to make room for icon)
        if heading:
            add_textbox(
                slide, _heading_truncate(heading, 8),
                left=left + PAD + icon_size + Inches(0.12),
                top=CT + Inches(0.12),
                width=col_width - 2 * PAD - icon_size - Inches(0.12),
                height=Inches(0.55),
                font_size=config.CARD_HEADING_SIZE, bold=True,
                color=config.COLOR_TEXT_DARK,
                font_name=config.TITLE_FONT,
            )

        # Bullet points below heading
        points = col.get("points", [])
        # Fallback: use generic message so column is never visually empty
        if not points:
            points = ["See document for details."]
        add_bullet_textbox(
            slide, [_truncate(p, 15) for p in points[:6]],
            left=left + PAD,
            top=CT + Inches(0.75),
            width=col_width - 2 * PAD,
            height=col_height - Inches(0.85),
            font_size=config.BODY_FONT_SIZE,
            color=config.COLOR_TEXT_DARK,
        )


# ---------------------------------------------------------------------------
# Layout: three_cards
# ---------------------------------------------------------------------------

def _three_cards(slide, data: dict) -> None:
    """Render a three-cards layout.

    Blueprint key expected:
        cards  (list of {number, heading, points})
    """
    cards = data.get("cards", [])[:3]
    if not cards:
        _single_focus(slide, data)
        return

    card_width = (CW - 2 * CGAP) / 3
    card_height = CH - GAP

    for i, card_data in enumerate(cards):
        left = ML + i * (card_width + CGAP)

        # Cream card with red border
        card = slide.shapes.add_shape(1, left, CT, card_width, card_height)
        style_shape(card, fill_color=config.COLOR_CARD_BG, line_color=config.COLOR_CARD_BORDER)

        # Red accent bar at top
        accent = slide.shapes.add_shape(1, left, CT, card_width, Inches(0.07))
        style_shape(accent, fill_color=config.COLOR_PRIMARY, line_color=None)

        # Large number badge
        number = card_data.get("number", str(i + 1).zfill(2))
        add_textbox(
            slide, number,
            left=left + PAD, top=CT + Inches(0.12),
            width=card_width - 2 * PAD, height=Inches(0.55),
            font_size=Pt(28), bold=True,
            color=config.COLOR_PRIMARY,
            font_name=config.TITLE_FONT,
        )

        # Small red icon in the top-right corner (heuristic from heading text).
        # Renders as a filled MSO_SHAPE; no external image required.
        icon_name = card_data.get("icon") or icon_for_text(
            card_data.get("heading", "")
        )
        icon_size = Inches(0.4)
        draw_icon_glyph(
            slide, icon_name,
            left=left + card_width - PAD - icon_size,
            top=CT + Inches(0.18),
            size=icon_size,
        )

        # Card heading below number
        heading = card_data.get("heading", "")
        if heading:
            add_textbox(
                slide, _heading_truncate(heading, 7),
                left=left + PAD, top=CT + Inches(0.72),
                width=card_width - 2 * PAD, height=Inches(0.55),
                font_size=config.CARD_HEADING_SIZE, bold=True,
                color=config.COLOR_TEXT_DARK,
                font_name=config.TITLE_FONT,
            )

        # Bullet points
        points = card_data.get("points", [])
        if not points:
            points = ["See document for details."]
        add_bullet_textbox(
            slide, [_truncate(p, 15) for p in points[:5]],
            left=left + PAD,
            top=CT + Inches(1.32),
            width=card_width - 2 * PAD,
            height=card_height - Inches(1.45),
            font_size=Pt(12),
            color=config.COLOR_TEXT_DARK,
        )


# ---------------------------------------------------------------------------
# Layout: key_stats
# ---------------------------------------------------------------------------

def _key_stats(slide, data: dict) -> None:
    """Render 2–4 big stat numbers with descriptive labels.

    Blueprint key expected:
        stats  (list of {value, label})
    """
    stats = data.get("stats", [])[:4]
    if not stats:
        _single_focus(slide, data)
        return

    n = len(stats)
    stat_width = (CW - (n - 1) * CGAP) / n
    stat_top = CT + Inches(0.3)
    stat_height = CH - Inches(0.6)

    for i, stat in enumerate(stats):
        left = ML + i * (stat_width + CGAP)

        # Cream card with red border
        box = slide.shapes.add_shape(1, left, stat_top, stat_width, stat_height)
        style_shape(box, fill_color=config.COLOR_CARD_BG, line_color=config.COLOR_CARD_BORDER)

        # Red accent bar at top
        accent = slide.shapes.add_shape(1, left, stat_top, stat_width, Inches(0.07))
        style_shape(accent, fill_color=config.COLOR_PRIMARY, line_color=None)

        # Vertically center the stat number + label within the card.
        # Guide rule: "Numbers not visually dominant enough" + "Not distributed properly"
        # Total content block = number (1.4") + gap (0.2") + label (0.6") = ~2.2"
        content_block = Inches(2.2)
        center_offset = (stat_height - Inches(0.07) - content_block) / 2  # subtract accent bar
        num_top = stat_top + Inches(0.07) + center_offset

        # Big number — centered and visually dominant (52pt bold)
        value = str(stat.get("value", ""))
        add_textbox(
            slide, value,
            left=left + PAD, top=num_top,
            width=stat_width - 2 * PAD, height=Inches(1.4),
            font_size=config.STAT_NUMBER_SIZE, bold=True,
            color=config.COLOR_PRIMARY,
            align=PP_ALIGN.CENTER,
            font_name=config.TITLE_FONT,
        )

        # Label below number — uses secondary color for better readability
        label = _truncate(str(stat.get("label", "")), 5)
        add_textbox(
            slide, label,
            left=left + PAD, top=num_top + Inches(1.6),
            width=stat_width - 2 * PAD, height=Inches(0.6),
            font_size=config.STAT_LABEL_SIZE,
            color=config.COLOR_TEXT_SECONDARY,
            align=PP_ALIGN.CENTER,
        )


# ---------------------------------------------------------------------------
# Layout: timeline
# ---------------------------------------------------------------------------

def _timeline(slide, data: dict) -> None:
    """Render a horizontal timeline with numbered steps.

    Blueprint key expected:
        steps  (list of {number, heading, description})
    """
    steps = data.get("steps", [])[:5]
    if not steps:
        _single_focus(slide, data)
        return

    n = len(steps)
    step_width = (CW - (n - 1) * CGAP) / n
    circle_r = Inches(0.3)
    line_y = CT + Inches(1.2)
    circle_top = line_y - circle_r
    circle_size = circle_r * 2

    # Horizontal connector line
    line = slide.shapes.add_shape(
        1,
        ML, line_y - Inches(0.02),
        CW, Inches(0.04),
    )
    style_shape(line, fill_color=config.COLOR_PRIMARY, line_color=None)

    for i, step in enumerate(steps):
        cx = ML + i * (step_width + CGAP) + step_width / 2

        # Red circle node
        circle = slide.shapes.add_shape(
            9,
            cx - circle_r, circle_top,
            circle_size, circle_size,
        )
        style_shape(circle, fill_color=config.COLOR_PRIMARY, line_color=None)

        # Step number in white inside circle
        number = step.get("number", str(i + 1).zfill(2))
        add_textbox(
            slide, number,
            left=cx - circle_r, top=circle_top,
            width=circle_size, height=circle_size,
            font_size=Pt(11), bold=True,
            color=pick_contrasting_text(config.COLOR_PRIMARY),
            align=PP_ALIGN.CENTER,
        )

        # Heading below circle
        heading = _heading_truncate(step.get("heading", ""), 6)
        add_textbox(
            slide, heading,
            left=ML + i * (step_width + CGAP),
            top=line_y + circle_r + Inches(0.12),
            width=step_width, height=Inches(0.5),
            font_size=Pt(13), bold=True,
            color=config.COLOR_TEXT_DARK,
            align=PP_ALIGN.CENTER,
        )

        # Description below heading — uses secondary color for better contrast
        desc = _truncate(step.get("description", ""), 15)
        add_textbox(
            slide, desc,
            left=ML + i * (step_width + CGAP),
            top=line_y + circle_r + Inches(0.68),
            width=step_width, height=Inches(1.8),
            font_size=Pt(11),
            color=config.COLOR_TEXT_SECONDARY,
            align=PP_ALIGN.CENTER,
        )


# ---------------------------------------------------------------------------
# Layout: process_flow
# ---------------------------------------------------------------------------

def _process_flow(slide, data: dict) -> None:
    """Render a left-to-right process flow with step boxes and arrow connectors.

    Blueprint key expected:
        steps  (list of {number, heading, description})
    """
    steps = data.get("steps", [])[:5]
    if not steps:
        _single_focus(slide, data)
        return

    n = len(steps)
    arrow_w = Inches(0.3)
    total_arrow = (n - 1) * arrow_w
    box_width = (CW - total_arrow) / n
    box_height = Inches(2.4)
    box_top = CT + (CH - box_height) / 2

    for i, step in enumerate(steps):
        left = ML + i * (box_width + arrow_w)

        box = slide.shapes.add_shape(1, left, box_top, box_width, box_height)
        style_shape(box, fill_color=config.COLOR_CARD_BG, line_color=config.COLOR_PRIMARY)

        number = step.get("number", str(i + 1).zfill(2))
        add_textbox(
            slide, number,
            left=left + PAD, top=box_top + PAD,
            width=box_width - 2 * PAD, height=Inches(0.45),
            font_size=Pt(20), bold=True, color=config.COLOR_PRIMARY,
            align=PP_ALIGN.CENTER,
        )

        heading = _heading_truncate(step.get("heading", ""), 6)
        add_textbox(
            slide, heading,
            left=left + PAD, top=box_top + Inches(0.6),
            width=box_width - 2 * PAD, height=Inches(0.55),
            font_size=Pt(13), bold=True, color=config.COLOR_TEXT_DARK,
            align=PP_ALIGN.CENTER,
        )

        desc = _truncate(step.get("description", ""), 15)
        add_textbox(
            slide, desc,
            left=left + PAD, top=box_top + Inches(1.2),
            width=box_width - 2 * PAD, height=Inches(1.1),
            font_size=Pt(11), color=config.COLOR_TEXT_SECONDARY,
            align=PP_ALIGN.CENTER,
        )

        if i < n - 1:
            arrow_left = left + box_width
            arrow = slide.shapes.add_shape(
                1,
                arrow_left,
                box_top + box_height / 2 - Inches(0.04),
                arrow_w, Inches(0.08),
            )
            style_shape(arrow, fill_color=config.COLOR_PRIMARY, line_color=None)


# ---------------------------------------------------------------------------
# Layout: comparison
# ---------------------------------------------------------------------------

def _comparison(slide, data: dict) -> None:
    """Render a high-contrast side-by-side comparison with two full-color columns.

    Blueprint keys expected:
        left.heading, left.points, right.heading, right.points
    """
    col_width = (CW - CGAP) / 2
    col_height = CH - GAP
    colors = [config.COLOR_PRIMARY, config.COLOR_TEXT_DARK]

    for i, side in enumerate(("left", "right")):
        col = data.get(side, {})
        if not isinstance(col, dict):
            col = {}

        left = ML + i * (col_width + CGAP)
        bg = colors[i]

        box = slide.shapes.add_shape(1, left, CT, col_width, col_height)
        style_shape(box, fill_color=bg, line_color=None)

        text_color = pick_contrasting_text(bg)

        # Icon glyph in top-left corner (white on colored background)
        heading = col.get("heading", "")
        icon_name = col.get("icon") or icon_for_text(heading)
        icon_size = Inches(0.4)
        draw_icon_glyph(
            slide, icon_name,
            left=left + PAD, top=CT + PAD,
            size=icon_size,
            fill=text_color,
        )

        if heading:
            add_textbox(
                slide, _heading_truncate(heading, 8),
                left=left + PAD, top=CT + PAD + icon_size + Inches(0.08),
                width=col_width - 2 * PAD, height=Inches(0.55),
                font_size=config.CARD_HEADING_SIZE, bold=True,
                color=text_color, align=PP_ALIGN.CENTER,
                font_name=config.TITLE_FONT,
            )

        points = col.get("points", [])
        if not points:
            points = ["See document for details."]
        add_bullet_textbox(
            slide, [_truncate(p, 15) for p in points[:6]],
            left=left + PAD, top=CT + PAD + icon_size + Inches(0.7),
            width=col_width - 2 * PAD,
            height=col_height - PAD - icon_size - Inches(0.8),
            font_size=config.BODY_FONT_SIZE, color=text_color,
        )


# ---------------------------------------------------------------------------
# Layout: icon_list
# ---------------------------------------------------------------------------

def _icon_list(slide, data: dict) -> None:
    """Render a numbered icon-circle list (3–4 rows).

    Blueprint key expected:
        items  (list of {number, heading, description})
    """
    items = data.get("items", [])[:4]
    if not items:
        _single_focus(slide, data)
        return

    row_height = min(CH / len(items) - CGAP, Inches(1.3))
    circle_size = Inches(0.55)
    text_left = ML + circle_size + Inches(0.25)
    text_width = CW - circle_size - Inches(0.25)

    for i, item in enumerate(items):
        row_top = CT + i * (row_height + CGAP / 2)
        circle_top = row_top + (row_height - circle_size) / 2

        # Red circle background for contrast
        circle = slide.shapes.add_shape(9, ML, circle_top, circle_size, circle_size)
        style_shape(circle, fill_color=config.COLOR_PRIMARY, line_color=None)

        # White icon glyph inside the circle — the glyph is a filled MSO_SHAPE
        # sized smaller than the circle so the red ring reads as a bezel.
        icon_name = item.get("icon") or icon_for_text(item.get("heading", ""))
        glyph_size = Inches(0.3)
        glyph_left = ML + (circle_size - glyph_size) / 2
        glyph_top = circle_top + (circle_size - glyph_size) / 2
        draw_icon_glyph(
            slide, icon_name,
            left=glyph_left, top=glyph_top,
            size=glyph_size,
            fill=pick_contrasting_text(config.COLOR_PRIMARY),
        )

        # Heading to the right
        heading = _heading_truncate(item.get("heading", ""), 8)
        add_textbox(
            slide, heading,
            left=text_left, top=row_top,
            width=text_width, height=Inches(0.4),
            font_size=config.CARD_HEADING_SIZE, bold=True,
            color=config.COLOR_TEXT_DARK,
            font_name=config.TITLE_FONT,
        )

        # Description below heading — secondary color for readable contrast
        desc = _truncate(item.get("description", ""), 20)
        add_textbox(
            slide, desc,
            left=text_left, top=row_top + Inches(0.42),
            width=text_width, height=row_height - Inches(0.45),
            font_size=config.BODY_FONT_SIZE,
            color=config.COLOR_TEXT_SECONDARY,
        )

        # Separator line (skip last row)
        if i < len(items) - 1:
            line = slide.shapes.add_shape(
                1,
                ML, row_top + row_height + CGAP / 4,
                CW, Inches(0.01),
            )
            style_shape(line, fill_color=config.COLOR_CARD_BORDER, line_color=None)


# ---------------------------------------------------------------------------
# Layout: single_focus
# ---------------------------------------------------------------------------

def _single_focus(slide, data: dict) -> None:
    """Render a single-focus slide: one large statement + supporting bullet list.

    Blueprint keys expected:
        focus   (str)  — the main key message displayed large
        points  (list[str]) — supporting bullet points below
    """
    focus = data.get("focus", "")
    points = data.get("points", [])

    if focus:
        # Large accent icon beside the focus statement
        icon_name = icon_for_text(focus)
        icon_size = Inches(0.5)
        draw_icon_glyph(
            slide, icon_name,
            left=ML, top=CT + Inches(0.05),
            size=icon_size,
        )

        add_textbox(
            slide, _truncate(focus, 20),
            left=ML + icon_size + Inches(0.15), top=CT,
            width=CW - icon_size - Inches(0.15), height=Inches(1.2),
            font_size=Pt(22), bold=True,
            color=config.COLOR_PRIMARY,
            align=PP_ALIGN.LEFT,
            font_name=config.TITLE_FONT,
        )

        # Thin red accent line separator
        line = slide.shapes.add_shape(
            1, ML, CT + Inches(1.25), Inches(1.5), Inches(0.05)
        )
        style_shape(line, fill_color=config.COLOR_PRIMARY, line_color=None)

        bullet_top = CT + Inches(1.4)
    else:
        bullet_top = CT

    if points:
        add_bullet_textbox(
            slide, [_truncate(p, 15) for p in points[:6]],
            left=ML, top=bullet_top,
            width=CW,
            height=CH - (bullet_top - CT) - GAP,
            font_size=config.BODY_FONT_SIZE,
            color=config.COLOR_TEXT_DARK,
        )


# ---------------------------------------------------------------------------
# Layout: six_cards
# A 3×2 grid of small icon-cards, each with heading + red divider + body.
# Designed for sections with 5–6 insights that don't need deep detail per
# card — an information-dense page that matches target page 3.
# ---------------------------------------------------------------------------

def _six_cards(slide, data: dict) -> None:
    """Render a 3×2 grid of small icon cards.

    Blueprint key expected:
        cards  (list of {heading, points|description})  — 4–6 items
    Falls back to deriving cards from left.points + right.points or flat points[].
    """
    # Collect cards from the various data shapes the blueprint might provide
    cards = data.get("cards", [])

    if not cards:
        # Harvest from left/right columns or flat points
        merged: list[str] = []
        for side in ("left", "right"):
            col = data.get(side, {})
            if isinstance(col, dict):
                merged.extend(col.get("points", []))
        if not merged:
            merged = data.get("points", [])
        for p in merged[:6]:
            words = p.split()
            heading = " ".join(words[:4]).rstrip(".,;:")
            cards.append({"heading": heading, "description": p})

    if not cards:
        _single_focus(slide, data)
        return

    cards = cards[:6]
    # Pad to at least 4 for visual balance
    while len(cards) < 4:
        cards.append({"heading": "Details", "description": "See document for details."})

    cols = 3
    rows = 2 if len(cards) > 3 else 1
    card_gap = Inches(0.25)
    card_w = (CW - (cols - 1) * card_gap) / cols
    card_h = (CH - GAP - (rows - 1) * card_gap) / rows

    for i, card_data in enumerate(cards):
        row = i // cols
        col = i % cols
        cx = ML + col * (card_w + card_gap)
        cy = CT + row * (card_h + card_gap)

        heading = card_data.get("heading", "")
        # Support both "description" (single string) and "points" (list)
        body = card_data.get("description", "")
        if not body:
            pts = card_data.get("points", [])
            body = ". ".join(pts[:2]) if pts else ""

        icon = card_data.get("icon") or icon_for_text(heading)

        draw_card_with_divider(
            slide,
            left=cx, top=cy,
            width=card_w, height=card_h,
            heading=_heading_truncate(heading, 5),
            body=_truncate(body, 22),
            icon_name=icon,
            heading_size=Pt(13),
            body_size=Pt(10),
        )


# ---------------------------------------------------------------------------
# Layout: five_cards_row
# Single horizontal row of 5 cards on a red full-bleed background.
# Matches target pages 2 and 10. Used for recap/conclusion slides
# where you want a punchy overview of 5 key items.
# This is a full-bleed layout — it paints its own title via a red top
# pill and skips add_slide_title in the dispatcher.
# ---------------------------------------------------------------------------

def _five_cards_row(slide, data: dict) -> None:
    """Render a horizontal row of 4–5 white icon-cards on a red background.

    Blueprint key expected:
        cards  (list of {heading, points|description})  — 4–5 items
    Also accepts: points (list[str]), or left/right column data.
    """
    title = data.get("title", "")

    # Red full-bleed background
    draw_red_full_background(slide)

    # Red top pill with the title (white text on slightly darker red pill)
    draw_red_top_pill(slide, title)

    # Collect card content
    cards = data.get("cards", [])

    if not cards:
        merged: list[str] = []
        for side in ("left", "right"):
            col = data.get(side, {})
            if isinstance(col, dict):
                merged.extend(col.get("points", []))
        if not merged:
            merged = data.get("points", [])
        # Also pull from focus + points for conclusion slides
        focus = data.get("focus", "")
        if focus and focus not in merged:
            merged.insert(0, focus)
        for p in merged[:5]:
            words = p.split()
            heading = " ".join(words[:4]).rstrip(".,;:")
            cards.append({"heading": heading, "description": p})

    if not cards:
        return

    cards = cards[:5]
    # Pad to at least 3 for visual balance
    while len(cards) < 3:
        cards.append({"heading": "Details", "description": "See document for details."})

    n = len(cards)
    card_gap = Inches(0.2)
    top = Inches(1.6)
    card_w = (CW - (n - 1) * card_gap) / n
    card_h = config.SLIDE_HEIGHT - top - Inches(0.5)

    for i, card_data in enumerate(cards):
        cx = ML + i * (card_w + card_gap)

        heading = card_data.get("heading", "")
        body = card_data.get("description", "")
        if not body:
            pts = card_data.get("points", [])
            body = ". ".join(pts[:2]) if pts else ""

        icon = card_data.get("icon") or icon_for_text(heading)

        draw_card_with_divider(
            slide,
            left=cx, top=top,
            width=card_w, height=card_h,
            heading=_heading_truncate(heading, 5),
            body=_truncate(body, 20),
            icon_name=icon,
            heading_size=Pt(13),
            body_size=Pt(10),
        )


# ---------------------------------------------------------------------------
# Layout: two_col_sidebar
# Red left sidebar (~32%) with section number + title. Content on the right
# is drawn from left.points / right.points as two stacked card sections,
# or from a flat points[] list split in half. This is the default layout
# for content-heavy sections and replaces the plain two_column for most
# middle slides. Accepts the same blueprint data shape as two_column.
# ---------------------------------------------------------------------------

def _two_col_sidebar(slide, data: dict) -> None:
    """Render a two-column layout with a red left sidebar.

    The sidebar holds the section title + optional section number.
    The right area shows two stacked content groups (from left/right or
    flat points), each rendered as a cream card with a serif heading
    and bullet points.

    render_content_slide() skips add_slide_title for this layout
    because the sidebar paints its own title.

    Blueprint keys accepted (same as two_column):
        left.heading, left.points, right.heading, right.points
    OR:
        points (list[str]) — auto-split into two halves
    """
    title = data.get("title", "")
    section_num = data.get("section_number")  # optional "01", "02", …

    # --- Left sidebar (reuse the visuals.py primitive) ---
    sidebar_right_edge = draw_red_left_sidebar(
        slide, title, section_number=section_num,
    )

    # --- Right content area ---
    right_left = sidebar_right_edge + Inches(0.4)
    right_top = Inches(0.5)
    right_width = config.SLIDE_WIDTH - right_left - Inches(0.5)
    right_bottom = config.SLIDE_HEIGHT - Inches(0.5)

    # Collect the two content groups
    left_col = data.get("left", {})
    right_col = data.get("right", {})
    if not isinstance(left_col, dict):
        left_col = {}
    if not isinstance(right_col, dict):
        right_col = {}

    # Fallback: if left/right are empty, split flat "points" list
    if not left_col.get("points") and not right_col.get("points"):
        flat = data.get("points", [])
        mid = max(1, len(flat) // 2)
        left_col = {"heading": "Overview", "points": flat[:mid]}
        right_col = {"heading": "Details", "points": flat[mid:]}

    groups = []
    if left_col.get("points"):
        groups.append(left_col)
    if right_col.get("points"):
        groups.append(right_col)
    if not groups:
        groups = [{"heading": "Overview", "points": ["See document for details."]}]

    # Render each group as a stacked card
    n = len(groups)
    card_gap = Inches(0.3)
    card_h = (right_bottom - right_top - (n - 1) * card_gap) / n

    for i, grp in enumerate(groups):
        card_top = right_top + i * (card_h + card_gap)
        heading = grp.get("heading", "")
        points = grp.get("points", [])
        if not points:
            points = ["See document for details."]

        # Cream card background with red border
        card = slide.shapes.add_shape(1, right_left, card_top, right_width, card_h)
        style_shape(card, fill_color=config.COLOR_CARD_BG, line_color=config.COLOR_CARD_BORDER)

        # Red accent stripe at top
        accent = slide.shapes.add_shape(1, right_left, card_top, right_width, Inches(0.06))
        style_shape(accent, fill_color=config.COLOR_PRIMARY, line_color=None)

        cursor = card_top + Inches(0.14)

        # Icon + heading row
        if heading:
            icon_name = icon_for_text(heading)
            icon_sz = Inches(0.35)
            draw_icon_glyph(
                slide, icon_name,
                left=right_left + PAD, top=cursor + Inches(0.02),
                size=icon_sz,
            )
            add_textbox(
                slide, _heading_truncate(heading, 8),
                left=right_left + PAD + icon_sz + Inches(0.15),
                top=cursor,
                width=right_width - 2 * PAD - icon_sz - Inches(0.15),
                height=Inches(0.45),
                font_size=config.CARD_HEADING_SIZE, bold=True,
                color=config.COLOR_TEXT_DARK,
                font_name=config.TITLE_FONT,
            )
            cursor += Inches(0.55)

        # Bullet points fill the rest of the card
        bullet_h = card_top + card_h - cursor - PAD
        if bullet_h > Inches(0.3):
            add_bullet_textbox(
                slide, [_truncate(p, 18) for p in points[:6]],
                left=right_left + PAD,
                top=cursor,
                width=right_width - 2 * PAD,
                height=bullet_h,
                font_size=config.BODY_FONT_SIZE,
                color=config.COLOR_TEXT_DARK,
            )


# ---------------------------------------------------------------------------
# Layout: exec_summary_with_photo
# A signature opening slide: a full-height red sidebar on the left holding
# a rotated "EXECUTIVE SUMMARY" badge, and a 2x2 grid of icon cards on the
# right pulling from the existing executive_summary schema (left.points +
# right.points). No photo is embedded; the "photo" in the name refers to
# the visual weight of the red sidebar matching the target deck.
# ---------------------------------------------------------------------------

def _exec_summary_with_photo(slide, data: dict) -> None:
    """Render the signature executive summary layout.

    Accepts either of the following blueprint shapes:
      1. Native: {"items": [{"heading": ..., "description": ...}, ...]}
      2. Legacy (from the executive_summary two_column schema):
         {"left": {"points": [...]}, "right": {"points": [...]}}
         — in which case the combined points are turned into 4 cards.

    The layout DOES NOT call add_slide_title — it paints its own sidebar
    title instead. render_content_slide() skips the normal title for this
    layout via the _FULL_BLEED_LAYOUTS set.
    """
    title = data.get("title", "Executive Summary")

    # --- Left sidebar: full-height red rectangle with rotated serif title ---
    sidebar_w = Inches(4.2)
    sidebar = slide.shapes.add_shape(
        1,          # MSO_SHAPE.RECTANGLE
        0, 0, sidebar_w, config.SLIDE_HEIGHT,
    )
    style_shape(sidebar, fill_color=config.COLOR_PRIMARY, line_color=None)
    sidebar_text_color = pick_contrasting_text(config.COLOR_PRIMARY)

    # Large centred "01" marker near the top of the sidebar (mirrors the
    # numbered rhythm from the target deck's section dividers).
    add_textbox(
        slide, "01",
        left=Inches(0.5), top=Inches(0.5),
        width=Inches(3.2), height=Inches(1.0),
        font_size=Pt(48), bold=True,
        color=sidebar_text_color,
        font_name=config.TITLE_FONT,
    )

    # Thin white divider line under the 01
    divider = slide.shapes.add_shape(1, Inches(0.5), Inches(1.55),
                                     Inches(0.8), Inches(0.03))
    style_shape(divider, fill_color=sidebar_text_color, line_color=None)

    # The main title — large, white, serif, left-aligned inside the sidebar
    add_textbox(
        slide, title.upper(),
        left=Inches(0.5), top=Inches(1.9),
        width=Inches(3.2), height=Inches(3.0),
        font_size=Pt(40), bold=True,
        color=sidebar_text_color,
        font_name=config.TITLE_FONT,
    )

    # Small eyebrow label at the bottom of the sidebar for balance
    add_textbox(
        slide, "KEY HIGHLIGHTS",
        left=Inches(0.5), top=Inches(6.4),
        width=Inches(3.2), height=Inches(0.4),
        font_size=Pt(10), bold=True,
        color=(
            RGBColor_light_pink()
            if sidebar_text_color == config.COLOR_TEXT_LIGHT
            else config.COLOR_TEXT_SECONDARY
        ),
    )

    # --- Right side: 2x2 grid of icon cards ---
    # Collect card content — prefer explicit "items", otherwise harvest
    # from left.points + right.points (the legacy two_column schema).
    items = data.get("items") or []
    if not items:
        merged: list[str] = []
        for side in ("left", "right"):
            col = data.get(side, {})
            if isinstance(col, dict):
                merged.extend(col.get("points", []))
        # Turn flat bullets into {heading, description} pairs
        for p in merged[:4]:
            words = p.split()
            heading = " ".join(words[:4]).rstrip(".,;:")
            items.append({"heading": heading, "description": p})

    if not items:
        return   # nothing to render on the right side

    # Fill up to 4 cards; 2x2 grid
    items = items[:4]
    grid_left = sidebar_w + Inches(0.5)
    grid_top = Inches(0.7)
    grid_width = config.SLIDE_WIDTH - grid_left - Inches(0.5)
    grid_height = config.SLIDE_HEIGHT - grid_top - Inches(0.6)

    cols = 2
    rows = 2 if len(items) > 2 else 1
    cell_w = (grid_width - Inches(0.3)) / cols
    cell_h = (grid_height - Inches(0.3)) / rows

    for i, item in enumerate(items):
        row = i // cols
        col = i % cols
        cx = grid_left + col * (cell_w + Inches(0.3))
        cy = grid_top + row * (cell_h + Inches(0.3))

        heading = item.get("heading", "")
        description = item.get("description", "")
        icon = item.get("icon") or icon_for_text(heading)

        draw_card_with_divider(
            slide,
            left=cx, top=cy,
            width=cell_w, height=cell_h,
            heading=_heading_truncate(heading, 6),
            body=_truncate(description, 28),
            icon_name=icon,
            heading_size=Pt(15),
            body_size=Pt(11),
        )


def RGBColor_light_pink():
    """Return a light pink color used for muted eyebrow labels on red sidebars.

    Wrapped as a function so the import line at the top of this module
    does not need to import RGBColor (layouts.py otherwise never touches
    colors directly).
    """
    from pptx.dml.color import RGBColor
    return RGBColor(0xFF, 0xD5, 0xCE)
