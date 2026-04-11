"""
renderer/tables.py — Styled table generation using python-pptx.

Tables are rendered with:
  - A bold, theme-colored header row (red background, white text)
  - Alternating light/white row backgrounds for readability
  - Consistent font sizes, padding, and column widths

Note: python-pptx's high-level cell fill API is limited, so cell
background colors are applied via direct XML manipulation.
"""

from __future__ import annotations

import logging

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn      # converts "a:tag" → "{namespace}tag" for lxml
from lxml import etree            # for direct XML manipulation of cell properties

import config
from renderer.utils import add_slide_title, add_slide_number, add_textbox, pick_contrasting_text

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Public entry point (called by engine.py)
# ---------------------------------------------------------------------------

def render_table_slide(slide, slide_data: dict, slide_num: int) -> None:
    """Render a complete table slide: title + table + optional caption + slide number.

    If the table fails to render (bad data, empty headers, etc.) a
    placeholder error message is shown instead.

    Args:
        slide: python-pptx slide object (already added to the presentation).
        slide_data: Single slide dict from the blueprint JSON.
        slide_num: Slide number to display in the bottom-right corner.
    """
    from pptx.enum.text import PP_ALIGN

    # --- Title ---
    title = slide_data.get("title", "")
    add_slide_title(slide, title)

    # Pull the table data dict and optional caption from the blueprint
    table_data = slide_data.get("table", {})
    caption = slide_data.get("caption", "")

    # Shrink table height if a caption needs to fit below it
    caption_height = Inches(0.4) if caption else 0
    table_height = config.CONTENT_HEIGHT - caption_height - config.ELEMENT_GAP

    # --- Table ---
    try:
        _add_styled_table(
            slide,
            table_data,
            left=config.MARGIN_LEFT,
            top=config.CONTENT_TOP,
            width=config.CONTENT_WIDTH,
            height=table_height,
        )
    except Exception as exc:
        # Show a visible error message if rendering fails
        logger.warning("Table render failed: %s", exc)
        add_textbox(
            slide, f"[Table could not be rendered: {exc}]",
            left=config.MARGIN_LEFT, top=config.CONTENT_TOP,
            width=config.CONTENT_WIDTH, height=table_height,
            font_size=config.BODY_FONT_SIZE, color=config.COLOR_TEXT_MUTED,
        )

    # --- Optional caption below the table ---
    if caption:
        add_textbox(
            slide, caption,
            left=config.MARGIN_LEFT,
            top=config.CONTENT_TOP + table_height + config.ELEMENT_GAP,
            width=config.CONTENT_WIDTH,
            height=Inches(0.35),
            font_size=config.CAPTION_FONT_SIZE,
            color=config.COLOR_TEXT_MUTED,
            align=PP_ALIGN.CENTER,
        )

    # --- Slide number ---
    add_slide_number(slide, slide_num)


# ---------------------------------------------------------------------------
# Internal table builder
# ---------------------------------------------------------------------------

def _add_styled_table(
    slide,
    table_data: dict,
    left: int,
    top: int,
    width: int,
    height: int,
) -> None:
    """Build a styled table shape and add it to the slide.

    Row 0 = styled header row (red background, white bold text, centered).
    Rows 1+ = data rows with alternating white/light-gray backgrounds.

    Args:
        slide: python-pptx slide object.
        table_data: Dict with "headers" (list[str]) and "rows" (list[list[str]]).
        left, top, width, height: Position and size in EMU.
    """
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])

    if not headers:
        raise ValueError("Table has no headers.")   # can't build a table without column names

    num_cols = len(headers)
    num_rows = len(rows) + 1   # +1 accounts for the header row

    # Protect against tables with too many rows overflowing the slide
    if num_rows > 20:
        rows = rows[:19]      # keep only first 19 data rows
        num_rows = 20
        logger.debug("Table truncated to 20 rows to prevent overflow.")

    # python-pptx creates the table grid with all cells empty
    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        left, top, width, height,
    )
    # Strip the inherited template table style so our manual cell colors take effect.
    # Without this, the template's dark table style overrides every fill we set.
    _clear_table_style(table_shape)
    table = table_shape.table

    # Distribute available width evenly across all columns
    col_width = width // num_cols
    for col in table.columns:
        col.width = col_width

    # --- Header row (row index 0) ---
    header_text_color = pick_contrasting_text(config.COLOR_HEADER_BG)
    for col_idx, header_text in enumerate(headers):
        cell = table.cell(0, col_idx)    # row=0 is the header row
        cell.text = str(header_text)
        _style_cell(
            cell,
            font_size=Pt(13),
            bold=True,
            text_color=header_text_color,
            bg_color=config.COLOR_HEADER_BG,       # red background
            align=PP_ALIGN.CENTER,                 # centred in header
        )

    # --- Red accent line below header row ---
    # A thin red bar spanning the table width sits just under the header
    # for extra visual weight (matches target deck styling).
    _add_header_accent(slide, left, top, width, table)

    # --- Data rows (row indices 1 through num_rows-1) ---
    for row_idx, row_data in enumerate(rows):
        # Alternate background: odd rows get light gray, even rows stay white
        bg = config.COLOR_ALT_ROW if row_idx % 2 == 1 else config.COLOR_TEXT_LIGHT

        for col_idx in range(num_cols):
            cell = table.cell(row_idx + 1, col_idx)   # +1 because row 0 is the header

            # Use empty string for cells where the source row has fewer columns
            value = str(row_data[col_idx]) if col_idx < len(row_data) else ""
            # Truncate long cell values to prevent overflow (max ~30 words)
            words = value.split()
            if len(words) > 30:
                value = " ".join(words[:30]) + "…"
            cell.text = value

            # First column acts as a row label — bold red text for emphasis
            is_label_col = col_idx == 0
            _style_cell(
                cell,
                font_size=Pt(12),
                bold=is_label_col,
                text_color=config.COLOR_PRIMARY if is_label_col else config.COLOR_TEXT_DARK,
                bg_color=bg,
                align=PP_ALIGN.LEFT,
            )


def _clear_table_style(table_shape) -> None:
    """Null out the table's inherited theme style.

    When a table is added to a slide that uses a template, python-pptx copies
    the template's default table style (often dark) into the table's tblPr.
    Setting the styleId to the null GUID disables that style so our manual
    per-cell fills are not overridden.

    Args:
        table_shape: The GraphicFrame shape containing the table.
    """
    try:
        tbl = table_shape.table._tbl
        tblPr = tbl.find(qn("a:tblPr"))
        if tblPr is not None:
            style_id_el = tblPr.find(qn("a:tableStyleId"))
            if style_id_el is not None:
                # Null GUID = "No Table Style" — disables all inherited formatting
                style_id_el.text = "{00000000-0000-0000-0000-000000000000}"
    except Exception as exc:
        logger.debug("Could not clear table style: %s", exc)


def _style_cell(
    cell,
    font_size: Pt,
    bold: bool,
    text_color: RGBColor,
    bg_color: RGBColor,
    align: PP_ALIGN,
) -> None:
    """Apply font styling and background fill to a single table cell.

    python-pptx doesn't expose a direct cell.fill API that works reliably
    across PowerPoint versions, so we inject the fill via raw XML.

    Args:
        cell: A python-pptx _Cell object.
        font_size: Font size (Pt).
        bold: Whether text should be bold.
        text_color: RGB color for the text.
        bg_color: RGB color for the cell background fill.
        align: Paragraph text alignment.
    """
    # --- Background fill via XML ---
    # We find-or-create <a:tcPr>, strip any existing fill elements (solidFill,
    # noFill, gradFill) that may have been inherited, then insert our own
    # solidFill. Stripping first prevents multiple competing fill elements.
    try:
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()

        # Remove any fill elements that may already exist (inherited from style)
        for fill_tag in ("a:solidFill", "a:noFill", "a:gradFill", "a:pattFill", "a:blipFill"):
            existing = tcPr.find(qn(fill_tag))
            if existing is not None:
                tcPr.remove(existing)

        # Insert our solid fill
        solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
        srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
        # RGBColor is a (R, G, B) tuple in this version of python-pptx
        hex_color = "{:02X}{:02X}{:02X}".format(
            bg_color[0], bg_color[1], bg_color[2]
        )
        srgbClr.set("val", hex_color)
    except Exception as exc:
        logger.debug("Cell background fill failed: %s", exc)

    # --- Text formatting ---
    tf = cell.text_frame
    tf.word_wrap = True   # prevent text from overflowing into adjacent cells
    # Disable auto-size so text stays at the specified size and wraps instead
    # of shrinking to fit (which can make dense tables unreadable).
    try:
        from pptx.enum.text import MSO_AUTO_SIZE
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    for para in tf.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.size = font_size
            run.font.bold = bold
            run.font.color.rgb = text_color

    # --- Vertical centering ---
    # Sets anchor="ctr" on tcPr so text is vertically middle-aligned in the cell.
    # This prevents text from sitting at the top of tall cells.
    try:
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        tcPr.set("anchor", "ctr")
    except Exception:
        pass   # vertical alignment is cosmetic — safe to skip on failure

    # --- Cell inner padding ---
    # Sets marL/R/T/B (margin left/right/top/bottom) on the tcPr element
    try:
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        margin = Inches(0.05)   # small inner padding so text doesn't touch cell borders
        tcPr.set("marL", str(int(margin)))
        tcPr.set("marR", str(int(margin)))
        tcPr.set("marT", str(int(margin)))
        tcPr.set("marB", str(int(margin)))
    except Exception:
        pass   # padding is cosmetic — safe to skip on failure


def _add_header_accent(slide, left: int, top: int, width: int, table) -> None:
    """Draw a thin red accent line just below the header row of a table.

    This adds extra visual weight to the header/body boundary, matching
    the target deck's table styling convention.

    Args:
        slide: python-pptx slide object.
        left: Table left position (EMU).
        top: Table top position (EMU).
        width: Table width (EMU).
        table: The python-pptx Table object.
    """
    from pptx.enum.shapes import MSO_SHAPE
    from renderer.utils import style_shape

    try:
        # Get the header row height to position the accent line
        header_h = table.rows[0].height
        accent_top = top + header_h
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, accent_top,
            width, Inches(0.04),
        )
        style_shape(accent, fill_color=config.COLOR_PRIMARY, line_color=None)
    except Exception:
        pass   # accent line is cosmetic — never fatal
