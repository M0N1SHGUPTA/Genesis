"""
renderer/utils.py — Shared helpers used across all renderer modules.

Every renderer module imports from here instead of duplicating code.
Functions cover: text box creation, shape styling, layout lookup,
slide numbers, title bars, and template slide removal.
"""

from __future__ import annotations

import logging
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn          # converts "a:tag" → "{namespace}tag" for lxml
from pptx.util import Inches, Pt, Emu
from lxml import etree               # for direct XML manipulation (corner radius, cell fills)

import config

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Layout lookup helpers
# ---------------------------------------------------------------------------

def get_layout_by_name(prs: Presentation, name_contains: str):
    """Find a slide layout whose name contains the given string (case-insensitive).

    Templates may have layouts in different order or with slightly different
    names across providers. This avoids hardcoding layout indices.

    Args:
        prs: The loaded Presentation object.
        name_contains: Substring to search for in layout names (e.g. "cover", "blank").

    Returns:
        The first matching SlideLayout, or layout index 2 (usually Blank) as fallback.
    """
    for layout in prs.slide_layouts:
        if name_contains.lower() in layout.name.lower():
            return layout   # found a match

    # No match found — warn and return a safe fallback
    logger.warning("Layout containing '%s' not found — falling back to index 2.", name_contains)
    idx = min(2, len(prs.slide_layouts) - 1)   # clamp in case template has fewer than 3 layouts
    return prs.slide_layouts[idx]


def get_blank_layout(prs: Presentation):
    """Return the Blank layout, trying common multilingual names before index fallback.

    "blank" = English, "leer" = German, "vide" = French.
    This ensures compatibility with templates localized for different languages.

    Args:
        prs: The loaded Presentation object.

    Returns:
        The Blank SlideLayout.
    """
    for name in ("blank", "leer", "vide"):
        for layout in prs.slide_layouts:
            if name in layout.name.lower():
                return layout
    # None of the known names matched — use index 2 as universal fallback
    return prs.slide_layouts[min(2, len(prs.slide_layouts) - 1)]


# ---------------------------------------------------------------------------
# Text box creation
# ---------------------------------------------------------------------------

def add_textbox(
    slide,
    text: str,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    font_size: Pt = config.BODY_FONT_SIZE,
    bold: bool = False,
    color: RGBColor = config.COLOR_TEXT_DARK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    word_wrap: bool = True,
    italic: bool = False,
) -> object:
    """Add a single-run styled text box to a slide.

    This is the core text placement function used everywhere in the renderer.
    All position/size values are in EMU (use Inches() or Pt() to convert).

    Args:
        slide: The python-pptx slide object to add the text box to.
        text: The string to display.
        left: Distance from the left edge of the slide (EMU).
        top: Distance from the top edge of the slide (EMU).
        width: Width of the text box (EMU).
        height: Height of the text box (EMU).
        font_size: Font size as a Pt value.
        bold: Whether the text should be bold.
        color: RGB text color.
        align: Paragraph alignment (LEFT, CENTER, RIGHT).
        word_wrap: If True, text wraps inside the box instead of overflowing.
        italic: Whether the text should be italic.

    Returns:
        The created TextBox shape object.
    """
    # Create the text box shape at the given position/size
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap   # prevents text from spilling outside the box boundary

    # python-pptx creates one empty paragraph by default — use it for the first run
    p = tf.paragraphs[0]
    p.alignment = align

    # A "run" is a contiguous span of text with the same formatting
    run = p.add_run()
    run.text = text

    # Apply font styling to this run
    font = run.font
    font.size = font_size
    font.bold = bold
    font.italic = italic
    font.color.rgb = color

    return txBox


def add_bullet_textbox(
    slide,
    points: list[str],
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    font_size: Pt = config.BODY_FONT_SIZE,
    color: RGBColor = config.COLOR_TEXT_DARK,
    bullet_char: str = "•",
) -> object:
    """Add a text box containing a bulleted list (max 6 items).

    Each bullet is a separate paragraph so spacing can be controlled.
    The bullet character is prepended manually to keep the approach simple
    and template-independent (no list XML required).

    Args:
        slide: The python-pptx slide object.
        points: List of bullet point strings. Capped at 6.
        left, top, width, height: Position/size in EMU.
        font_size: Font size for all bullet text.
        color: RGB text color for all bullet text.
        bullet_char: Character placed before each bullet (default "•").

    Returns:
        The created TextBox shape object.
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True   # essential — bullets can be long

    # Add one paragraph per bullet point (hard cap at 6 to prevent text overflow)
    for i, point in enumerate(points[:6]):
        # First paragraph already exists; add new ones for subsequent bullets
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)   # small vertical gap above each bullet for readability

        run = p.add_run()
        # Prefix the bullet character + two spaces for visual indent
        run.text = f"{bullet_char}  {point}"
        run.font.size = font_size
        run.font.color.rgb = color

    return txBox


# ---------------------------------------------------------------------------
# Shape styling
# ---------------------------------------------------------------------------

def style_shape(
    shape,
    fill_color: Optional[RGBColor] = None,
    line_color: Optional[RGBColor] = None,
    line_width: Pt = config.CARD_BORDER_WIDTH,
    corner_radius: Emu = config.CARD_CORNER_RADIUS,
) -> None:
    """Apply fill color, border color/width, and corner radius to a shape.

    Used for cards, background boxes, divider bars, and arrow connectors.

    Args:
        shape: A python-pptx shape object (rectangle, oval, etc.).
        fill_color: Solid background fill color. Pass None for transparent fill.
        line_color: Border/outline color. Pass None for no border.
        line_width: Border line thickness (Pt). Only used when line_color is set.
        corner_radius: Rounded corner radius in EMU. Only takes effect on
                       shapes whose geometry is "roundRect".
    """
    # --- Fill ---
    if fill_color is not None:
        shape.fill.solid()                      # switch fill type to solid
        shape.fill.fore_color.rgb = fill_color  # set the fill color
    else:
        shape.fill.background()   # transparent fill (shows slide background through)

    # --- Border/outline ---
    if line_color is not None:
        shape.line.color.rgb = line_color   # set border color
        shape.line.width = line_width       # set border thickness
    else:
        shape.line.fill.background()   # no border

    # --- Corner radius (only for roundRect geometry) ---
    # python-pptx doesn't expose corner radius directly, so we edit the XML.
    try:
        sp = shape._element                          # the underlying <p:sp> XML element
        spPr = sp.find(qn("p:spPr"))                 # shape properties element
        if spPr is not None:
            prstGeom = spPr.find(qn("a:prstGeom"))   # preset geometry element
            if prstGeom is not None and prstGeom.get("prst") == "roundRect":
                # The adjust value list controls the corner curve amount
                avLst = prstGeom.find(qn("a:avLst"))
                if avLst is None:
                    avLst = etree.SubElement(prstGeom, qn("a:avLst"))
                gd = avLst.find(qn("a:gd"))    # guide definition
                if gd is None:
                    gd = etree.SubElement(avLst, qn("a:gd"))
                gd.set("name", "adj")
                gd.set("fmla", f"val {int(corner_radius)}")   # set the radius value
    except Exception:
        pass   # shape doesn't support corner radius (e.g. oval) — silently skip


# ---------------------------------------------------------------------------
# Slide number footer
# ---------------------------------------------------------------------------

def add_slide_number(slide, slide_num: int) -> None:
    """Add a small slide number label to the bottom-right corner.

    Placed inside the safe margin zone so it doesn't get cut off.

    Args:
        slide: The python-pptx slide object.
        slide_num: Integer slide number to display.
    """
    # Calculate position: align to right margin, sit above bottom margin
    right_edge = config.SLIDE_WIDTH - config.MARGIN_RIGHT
    bottom_edge = config.SLIDE_HEIGHT - Inches(0.35)
    width = Inches(0.5)
    height = Inches(0.25)

    add_textbox(
        slide,
        str(slide_num),
        left=right_edge - width,    # flush against right margin
        top=bottom_edge - height,
        width=width,
        height=height,
        font_size=config.SLIDE_NUM_SIZE,
        color=config.COLOR_TEXT_MUTED,   # muted gray so it doesn't compete with content
        align=PP_ALIGN.RIGHT,
    )


# ---------------------------------------------------------------------------
# Slide title bar
# ---------------------------------------------------------------------------

def add_slide_title(
    slide,
    title: str,
    color: RGBColor = config.COLOR_TEXT_DARK,
) -> None:
    """Add a standardised bold title text box at the top of a content slide.

    All content slides use this to ensure consistent title positioning.
    The title occupies the area from MARGIN_TOP down to CONTENT_TOP.

    Args:
        slide: The python-pptx slide object.
        title: The slide title string.
        color: Text color (defaults to near-black).
    """
    add_textbox(
        slide,
        title,
        left=config.MARGIN_LEFT,
        top=config.MARGIN_TOP,
        width=config.CONTENT_WIDTH,
        height=Inches(1.2),              # tall enough for 2-line titles at 32pt
        font_size=config.TITLE_FONT_SIZE,
        bold=True,
        color=color,
        align=PP_ALIGN.LEFT,             # titles are always left-aligned
    )


# ---------------------------------------------------------------------------
# Template slide removal
# ---------------------------------------------------------------------------

def remove_template_slides(prs: Presentation, keep_from: int) -> None:
    """Remove the original template placeholder slides from the presentation.

    When we load a template .pptx, it already contains several placeholder
    slides (the template demos). All our new content slides are appended
    after them. This function deletes the originals so only our slides remain.

    Deletion is done via direct XML manipulation because python-pptx does not
    expose a public slide-delete API.

    Args:
        prs: The Presentation object.
        keep_from: Number of slides to remove from the front (i.e. the count
                   of original template slides recorded before we started adding).
    """
    xml_slides = prs.slides._sldIdLst   # the XML list of slide ID references

    for _ in range(keep_from):
        if len(xml_slides) == 0:
            break   # nothing left to remove

        # Each entry in _sldIdLst holds an r:id relationship reference
        slide_elem = xml_slides[0]
        rId = slide_elem.get(qn("r:id"))

        if rId:
            try:
                # Drop the relationship from the presentation part
                # (this removes the file reference inside the .pptx ZIP)
                prs.part.drop_rel(rId)
            except Exception as exc:
                logger.debug("drop_rel failed for rId %s: %s", rId, exc)

        # Remove the slide ID entry from the list
        del xml_slides[0]

    logger.debug("Removed %d template placeholder slide(s).", keep_from)
