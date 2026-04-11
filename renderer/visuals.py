"""
renderer/visuals.py — Reusable decorative primitives.

Higher-level building blocks that compose text + shapes into visual motifs
shared by multiple layouts. Every primitive here draws onto a python-pptx
slide using only programmatic shapes (no external images), so the output
stays editable in PowerPoint / Google Slides / LibreOffice.

What belongs here vs renderer/utils.py:
  - utils.py   → low-level helpers (add_textbox, style_shape, placeholders)
  - visuals.py → higher-level motifs (sidebars, pills, icon glyphs, cards
                 with dividers) that stack multiple shapes into one unit

Primitives:
  draw_red_left_sidebar      — red full-height rectangle with title text
  draw_red_top_pill          — red pill-shaped title bar at top of slide
  draw_red_full_background   — fill entire slide with primary red
  draw_numbered_badge        — filled red oval with a white number
  draw_card_with_divider     — white card with thin red divider under heading
  draw_icon_glyph            — MSO_SHAPE glyph filled in primary red
  apply_serif_font           — set a run's font to the project serif
  icon_for_text              — heuristic keyword → MSO_SHAPE mapping
"""

from __future__ import annotations

import logging
import re
from typing import Optional

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

import config
from renderer.utils import (
    add_textbox,
    pick_contrasting_text,
    strip_numeric_prefix,
    style_shape,
)

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Font helpers
# ---------------------------------------------------------------------------

# Serif stack used for titles and card headings. PowerPoint will fall through
# to the next available font if the first is not installed on the host.
SERIF_TITLE_FONT = "Georgia"


def apply_serif_font(run) -> None:
    """Set a python-pptx run's font.name to the project serif face.

    Called wherever we want a title or heading to read as editorial rather
    than default sans.

    Args:
        run: A python-pptx _Run object (returned by paragraph.add_run()).
    """
    try:
        run.font.name = SERIF_TITLE_FONT
    except Exception:
        pass   # font assignment failures are never fatal


# ---------------------------------------------------------------------------
# Red left sidebar
# ---------------------------------------------------------------------------

def draw_red_left_sidebar(
    slide,
    title: str,
    subtitle: str = "",
    *,
    width_frac: float = 0.32,
    section_number: Optional[str] = None,
) -> Emu:
    """Draw a full-height red sidebar on the left edge of the slide.

    Used by content layouts that want the target-deck "red column + content
    on the right" motif. The sidebar holds an optional large section number,
    a serif title, and an optional muted subtitle.

    Args:
        slide:          python-pptx slide object.
        title:          Sidebar title (rendered large, white, serif).
        subtitle:       Optional smaller subtitle below the title.
        width_frac:     Sidebar width as fraction of the slide width.
        section_number: Optional "01" / "02" style number shown above title.

    Returns:
        The Emu x-coordinate where the sidebar ends (useful as the left edge
        of whatever content should sit to its right).
    """
    title = strip_numeric_prefix(title)
    sidebar_w = Emu(int(config.SLIDE_WIDTH * width_frac))
    sidebar_text_color = pick_contrasting_text(config.COLOR_PRIMARY)
    subtitle_color = (
        RGBColor(0xFF, 0xE5, 0xE1)
        if sidebar_text_color == config.COLOR_TEXT_LIGHT
        else config.COLOR_TEXT_SECONDARY
    )

    # Full-height red rectangle
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, sidebar_w, config.SLIDE_HEIGHT,
    )
    style_shape(bar, fill_color=config.COLOR_PRIMARY, line_color=None)

    inner_left = Inches(0.5)
    inner_width = sidebar_w - Inches(1.0)

    # Optional big section number ("01", "02", …) above the title
    if section_number:
        add_textbox(
            slide, section_number,
            left=inner_left, top=Inches(1.0),
            width=inner_width, height=Inches(1.0),
            font_size=Pt(56), bold=True,
            color=sidebar_text_color,
        )

    # Title — serif, big, white
    title_top = Inches(2.1) if section_number else Inches(1.8)
    title_box = add_textbox(
        slide, title,
        left=inner_left, top=title_top,
        width=inner_width, height=Inches(2.8),
        font_size=Pt(28), bold=True,
        color=sidebar_text_color,
    )
    # Upgrade the title run to serif
    for para in title_box.text_frame.paragraphs:
        for run in para.runs:
            apply_serif_font(run)

    # Optional subtitle below title
    if subtitle:
        add_textbox(
            slide, subtitle,
            left=inner_left, top=title_top + Inches(2.9),
            width=inner_width, height=Inches(1.2),
            font_size=Pt(12),
            color=subtitle_color,
        )

    return sidebar_w


# ---------------------------------------------------------------------------
# Red top pill (rounded title bar)
# ---------------------------------------------------------------------------

def draw_red_top_pill(slide, title: str) -> None:
    """Draw a rounded red pill at the top of the slide holding a white title.

    Replaces the plain black title textbox for layouts that want more visual
    weight in the title bar (matches target pages 6 and 10).

    Args:
        slide: python-pptx slide object.
        title: Title text.
    """
    left = config.MARGIN_LEFT
    top = Inches(0.35)
    width = config.CONTENT_WIDTH
    height = Inches(0.9)

    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height,
    )
    # Push the corner radius toward max so it reads as a pill
    style_shape(
        pill,
        fill_color=config.COLOR_PRIMARY,
        line_color=None,
        corner_radius=Emu(50000),
    )

    title_box = add_textbox(
        slide, title,
        left=left + Inches(0.4), top=top + Inches(0.12),
        width=width - Inches(0.8), height=Inches(0.7),
        font_size=Pt(22), bold=True,
        color=pick_contrasting_text(config.COLOR_PRIMARY),
        align=PP_ALIGN.LEFT,
    )
    for para in title_box.text_frame.paragraphs:
        for run in para.runs:
            apply_serif_font(run)


# ---------------------------------------------------------------------------
# Full-bleed red background
# ---------------------------------------------------------------------------

def draw_red_full_background(slide) -> None:
    """Fill the entire slide with primary red.

    Used for layouts that put white cards on a red page. Drawn first so all
    subsequent shapes sit on top.

    Args:
        slide: python-pptx slide object.
    """
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, config.SLIDE_WIDTH, config.SLIDE_HEIGHT,
    )
    style_shape(bg, fill_color=config.COLOR_PRIMARY, line_color=None)


# ---------------------------------------------------------------------------
# Numbered badge (oval with number)
# ---------------------------------------------------------------------------

def draw_numbered_badge(
    slide,
    number: str,
    left: Emu,
    top: Emu,
    size: Emu = Inches(0.55),
    *,
    fill: RGBColor = None,
    text_color: RGBColor = None,
    font_size: Pt = Pt(14),
) -> None:
    """Draw a filled circle containing a small bold number.

    Used for card numbers (01/02/03), timeline nodes, and icon-list markers.

    Args:
        slide:      python-pptx slide object.
        number:     Text to display inside the badge (typically "01", "02", …).
        left, top:  Top-left position of the circle (EMU).
        size:       Diameter of the circle (EMU).
        fill:       Circle fill color. Defaults to primary red.
        text_color: Number color. Defaults to white.
        font_size:  Font size for the number.
    """
    if fill is None:
        fill = config.COLOR_PRIMARY
    if text_color is None:
        text_color = pick_contrasting_text(fill)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    style_shape(circle, fill_color=fill, line_color=None)

    # Centred number (textbox spans the same rect as the oval so the centring
    # calculation is just PP_ALIGN.CENTER + vertical centring of a 1-line run).
    add_textbox(
        slide, number,
        left=left, top=top,
        width=size, height=size,
        font_size=font_size, bold=True,
        color=text_color,
        align=PP_ALIGN.CENTER,
    )


# ---------------------------------------------------------------------------
# Card with thin red divider under heading
# ---------------------------------------------------------------------------

def draw_card_with_divider(
    slide,
    left: Emu,
    top: Emu,
    width: Emu,
    height: Emu,
    heading: str,
    body: str,
    *,
    icon_name: Optional[str] = None,
    heading_size: Pt = Pt(14),
    body_size: Pt = Pt(11),
) -> None:
    """Draw a white card with an icon, heading, thin red divider, then body.

    Matches the target deck's small-card style (target pages 3 and 10).
    Stacked elements inside the card:
        [ icon ]
        heading  (serif, bold)
        ─── (thin red divider line)
        body text (muted gray, 2 lines)

    Args:
        slide:     python-pptx slide object.
        left, top: Top-left of the card (EMU).
        width, height: Card dimensions (EMU).
        heading:   Card heading (bold, serif).
        body:      Card body (1–2 lines of detail).
        icon_name: Optional icon keyword; passed through draw_icon_glyph.
        heading_size: Font size for the heading.
        body_size:    Font size for the body.
    """
    pad = Inches(0.18)

    # White card with a thin red border
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    style_shape(
        card,
        fill_color=config.COLOR_TEXT_LIGHT,
        line_color=config.COLOR_CARD_BORDER,
    )

    cursor_top = top + pad

    # Optional icon (small red glyph)
    if icon_name:
        icon_size = Inches(0.35)
        draw_icon_glyph(
            slide, icon_name,
            left=left + pad, top=cursor_top,
            size=icon_size,
        )
        cursor_top = cursor_top + icon_size + Inches(0.08)

    # Heading
    heading_h = Inches(0.42)
    heading_box = add_textbox(
        slide, heading,
        left=left + pad, top=cursor_top,
        width=width - 2 * pad, height=heading_h,
        font_size=heading_size, bold=True,
        color=config.COLOR_TEXT_DARK,
    )
    for para in heading_box.text_frame.paragraphs:
        for run in para.runs:
            apply_serif_font(run)
    cursor_top = cursor_top + heading_h

    # Thin red divider line
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + pad, cursor_top + Inches(0.02),
        Inches(0.8), Inches(0.025),
    )
    style_shape(divider, fill_color=config.COLOR_PRIMARY, line_color=None)
    cursor_top = cursor_top + Inches(0.12)

    # Body text (remaining space)
    body_h = top + height - cursor_top - pad
    if body_h > Inches(0.25):
        add_textbox(
            slide, body,
            left=left + pad, top=cursor_top,
            width=width - 2 * pad, height=body_h,
            font_size=body_size,
            color=config.COLOR_TEXT_DARK,
        )


# ---------------------------------------------------------------------------
# Icon glyph — draws a filled MSO_SHAPE that reads as an icon
# ---------------------------------------------------------------------------

# Keyword → MSO_SHAPE mapping. The key list is intentionally small so that
# each icon is recognisable and readable at ~0.35" size on a slide.
# Insertion order matters: earlier keys win on ambiguous matches.
_ICON_MAP: list[tuple[str, int]] = [
    # (keyword regex, MSO_SHAPE int id)
    (r"\b(growth|increase|rise|up|expansion|expand)\b",     MSO_SHAPE.UP_ARROW),
    (r"\b(decline|decrease|fall|drop|down|loss)\b",         MSO_SHAPE.DOWN_ARROW),
    (r"\b(acqui|merger|m&a|deal|buyout|takeover|divest)\b", MSO_SHAPE.STAR_5_POINT),
    (r"\b(learn|educat|train|reskill|course|certif|"
     r"academ|school|university)\b",                        MSO_SHAPE.RECTANGLE),
    (r"\b(invest|capital|fund|budget|cost|price|revenue|"
     r"financ|money|trillion|billion|dollar)\b",            MSO_SHAPE.OVAL),
    (r"\b(ai|intelligen|tech|power|electric|energy|"
     r"innovation|innov|digital|automat|software|"
     r"machine|neural|comput|algorithm|agent)\b",           MSO_SHAPE.LIGHTNING_BOLT),
    (r"\b(secur|risk|threat|defens|protect|shield|"
     r"cyber|identity|privacy)\b",                          MSO_SHAPE.PENTAGON),
    (r"\b(consult|advisory|professional|service|"
     r"practice|deliver)\b",                                MSO_SHAPE.DIAMOND),
    (r"\b(engineer|manufactur|aerospace|industrial|"
     r"factory|infrastr|construct)\b",                      MSO_SHAPE.HEXAGON),
    (r"\b(people|talent|workforce|employ|human|"
     r"citizen|community|team|headcount)\b",                MSO_SHAPE.HEART),
    (r"\b(global|world|international|nation|region|"
     r"country|cross-border|geographic)\b",                 MSO_SHAPE.DONUT),
    (r"\b(time|speed|deadline|schedul|phase|timeline|"
     r"quarter|year|month)\b",                              MSO_SHAPE.RIGHT_ARROW),
    (r"\b(strategy|goal|target|objective|vision|"
     r"mission|plan|recommend|priorit)\b",                  MSO_SHAPE.STAR_5_POINT),
    (r"\b(integrat|consolidat|embed|align|unif|transform|"
     r"reinvent|restructur)\b",                             MSO_SHAPE.DONUT),
    (r"\b(compar|versus|contrast|alternative|option)\b",    MSO_SHAPE.DIAMOND),
    (r"\b(data|analytic|metric|measur|report|dashboard)\b", MSO_SHAPE.RECTANGLE),
    (r"\b(govern|polic|regulat|law|compliance)\b",          MSO_SHAPE.PENTAGON),
    (r"\b(climate|environment|green|sustain|carbon)\b",     MSO_SHAPE.SUN),
    (r"\b(insur|underwrit|claim|actuar|pension)\b",         MSO_SHAPE.PENTAGON),
    (r"\b(health|care|medic|well)\b",                       MSO_SHAPE.HEART),
    (r"\b(market|position|competi|benchmark|leader)\b",     MSO_SHAPE.UP_ARROW),
    (r"\b(revenue|profit|earning|margin|fiscal)\b",         MSO_SHAPE.OVAL),
]


def icon_for_text(text: str) -> str:
    """Heuristically map an arbitrary string to an MSO_SHAPE icon name.

    This runs at render-time on card/step headings so we do not need to
    teach the LLM anything about icons. If no keyword matches, returns
    a generic "dot" fallback.

    Args:
        text: Any string (card heading, step title, etc.).

    Returns:
        A string key like "up_arrow", "lightning_bolt", "dot", … suitable
        for passing to draw_icon_glyph.
    """
    if not text:
        return "dot"

    lowered = text.lower()
    for pattern, shape_id in _ICON_MAP:
        if re.search(pattern, lowered):
            # Map the MSO_SHAPE enum int back to its lowercase name for the API.
            return _MSO_ID_TO_NAME.get(shape_id, "dot")
    return "dot"


# Reverse lookup: MSO_SHAPE int id → canonical name string used by draw_icon_glyph
_MSO_ID_TO_NAME: dict[int, str] = {
    int(MSO_SHAPE.UP_ARROW):       "up_arrow",
    int(MSO_SHAPE.DOWN_ARROW):     "down_arrow",
    int(MSO_SHAPE.RIGHT_ARROW):    "right_arrow",
    int(MSO_SHAPE.LIGHTNING_BOLT): "lightning_bolt",
    int(MSO_SHAPE.PENTAGON):       "pentagon",
    int(MSO_SHAPE.HEART):          "heart",
    int(MSO_SHAPE.DONUT):          "donut",
    int(MSO_SHAPE.STAR_5_POINT):   "star",
    int(MSO_SHAPE.DIAMOND):        "diamond",
    int(MSO_SHAPE.SUN):            "sun",
    int(MSO_SHAPE.RECTANGLE):      "square",
    int(MSO_SHAPE.OVAL):           "circle",
    int(MSO_SHAPE.HEXAGON):        "hexagon",
}

# Forward lookup: name string → MSO_SHAPE enum (used by draw_icon_glyph)
_NAME_TO_MSO: dict[str, int] = {
    "up_arrow":       MSO_SHAPE.UP_ARROW,
    "down_arrow":     MSO_SHAPE.DOWN_ARROW,
    "right_arrow":    MSO_SHAPE.RIGHT_ARROW,
    "lightning_bolt": MSO_SHAPE.LIGHTNING_BOLT,
    "pentagon":       MSO_SHAPE.PENTAGON,
    "heart":          MSO_SHAPE.HEART,
    "donut":          MSO_SHAPE.DONUT,
    "star":           MSO_SHAPE.STAR_5_POINT,
    "diamond":        MSO_SHAPE.DIAMOND,
    "sun":            MSO_SHAPE.SUN,
    "square":         MSO_SHAPE.RECTANGLE,
    "circle":         MSO_SHAPE.OVAL,
    "hexagon":        MSO_SHAPE.HEXAGON,
    "dot":            MSO_SHAPE.OVAL,
}


def draw_icon_glyph(
    slide,
    icon_name: str,
    left: Emu,
    top: Emu,
    size: Emu = Inches(0.35),
    *,
    fill: RGBColor = None,
) -> None:
    """Draw a small filled MSO_SHAPE that reads as an icon.

    The shape is sized as a square bounding box and rendered in the primary
    red by default. Falls back to a small red circle ("dot") if the icon
    name is unknown.

    Args:
        slide:     python-pptx slide object.
        icon_name: String name from icon_for_text() or a direct key like "sun".
        left, top: Top-left of the icon bounding box (EMU).
        size:      Icon width and height (EMU). Icons are square.
        fill:      Fill color. Defaults to primary red.
    """
    if fill is None:
        fill = config.COLOR_PRIMARY

    shape_enum = _NAME_TO_MSO.get(icon_name, MSO_SHAPE.OVAL)
    try:
        shape = slide.shapes.add_shape(shape_enum, left, top, size, size)
        style_shape(shape, fill_color=fill, line_color=None)
    except Exception as exc:
        logger.debug("draw_icon_glyph failed for %s: %s — drawing dot.", icon_name, exc)
        # Absolute fallback — a small filled dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
        style_shape(dot, fill_color=fill, line_color=None)
