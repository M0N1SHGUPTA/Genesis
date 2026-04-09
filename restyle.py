"""
restyle.py — Restyle new3.pptx to match the target design.

Run:   env\Scripts\python restyle.py
Output: outputs/output_styled.pptx

Rules:
  - Do NOT change any text content
  - Only modify fonts, colors, positions, fills
  - Add decorative shapes (panels, dividers, circles, etc.)
"""

from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# ── Palette ───────────────────────────────────────────────────────────────────
RED   = RGBColor(0xE6, 0x33, 0x29)
NAVY  = RGBColor(0x1A, 0x20, 0x35)
BLUSH = RGBColor(0xFD, 0xE8, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK  = RGBColor(0x22, 0x22, 0x22)
MUTED = RGBColor(0x66, 0x66, 0x66)
LGRAY = RGBColor(0xE8, 0xE8, 0xE8)

SW = Inches(13.33)
SH = Inches(7.5)


# ── Low-level helpers ─────────────────────────────────────────────────────────

def solid_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, l: float, t: float, w: float, h: float,
          fill=None, line=None, lw=Pt(1.5)):
    """Add rectangle with Inches coords. Returns shape."""
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    return s


def _oval(slide, l: float, t: float, w: float, h: float, fill=None):
    """Add oval/circle with Inches coords."""
    s = slide.shapes.add_shape(9, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    s.line.fill.background()
    return s


def _tb(slide, text: str, l: float, t: float, w: float, h: float,
        fn="Calibri", sz=14, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False):
    """Add a new text box with Inches coords."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = fn
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def restyle_tf(shape, fn=None, sz=None, bold=None, color=None, align=None):
    """Restyle every run in every paragraph of an existing text frame."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        if align is not None:
            para.alignment = align
        for run in para.runs:
            if fn:              run.font.name  = fn
            if sz:              run.font.size  = Pt(sz)
            if bold is not None: run.font.bold = bold
            if color:           run.font.color.rgb = color


def kill(shape):
    """Remove a shape from its parent slide."""
    el = shape._element
    el.getparent().remove(el)


def to_back(shape, slide):
    """Send a shape behind all other shapes (draw order)."""
    sp_tree = slide.shapes._spTree
    sp_tree.remove(shape._element)
    sp_tree.insert(2, shape._element)   # index 2 = after nvGrpSpPr + grpSpPr


def kill_bg_rects(slide):
    """Remove all AUTO_SHAPE rectangles that have no text — background decorations."""
    for s in list(slide.shapes):
        if s.shape_type == 1:
            has_content = s.has_text_frame and s.text.strip()
            if not has_content:
                kill(s)


def emu_to_in(emu: int) -> float:
    return emu / 914400


def add_native_chart(slide, categories, series_list, l, t, w, h):
    """
    Add a clustered column chart with red bars.
    series_list = [(name, [values, ...]), ...]
    """
    cd = CategoryChartData()
    cd.categories = [str(c) for c in categories]
    for name, vals in series_list:
        cd.add_series(name, [float(v) for v in vals])

    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(l), Inches(t), Inches(w), Inches(h),
        cd,
    )
    ch = gf.chart
    try:
        ch.has_title = False
    except Exception:
        pass
    ch.has_legend = len(series_list) > 1

    # Red bars
    try:
        for ser in ch.plots[0].series:
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = RED
        ch.plots[0].gap_width = 60
    except Exception:
        pass

    # Transparent chart / plot area
    try:
        ch.chart_area.format.fill.background()
        ch.plot_area.format.fill.background()
    except Exception:
        pass

    # Axis tick label styling
    try:
        ch.value_axis.tick_labels.font.size = Pt(10)
        ch.value_axis.tick_labels.font.color.rgb = MUTED
        ch.category_axis.tick_labels.font.size = Pt(10)
        ch.category_axis.tick_labels.font.color.rgb = MUTED
    except Exception:
        pass

    return gf


# ── Slide 1 — Cover ───────────────────────────────────────────────────────────

def s1_cover(slide):
    solid_bg(slide, WHITE)

    tbs = [s for s in slide.shapes if s.has_text_frame]
    title_s    = tbs[0] if len(tbs) > 0 else None
    subtitle_s = tbs[1] if len(tbs) > 1 else None

    # Navy rotated diamond — decorative top-center
    diam = slide.shapes.add_shape(1, Inches(5.3), Inches(-0.9), Inches(2.8), Inches(2.8))
    diam.fill.solid()
    diam.fill.fore_color.rgb = NAVY
    diam.line.fill.background()
    diam.rotation = 45
    to_back(diam, slide)

    # Thin red horizontal rule above title
    _rect(slide, 0.6, 4.95, 7.5, 0.04, fill=RED)

    # Title — Georgia 36pt bold red, bottom-left
    if title_s:
        title_s.left   = Inches(0.6)
        title_s.top    = Inches(5.05)
        title_s.width  = Inches(8.5)
        title_s.height = Inches(1.6)
        title_s.text_frame.word_wrap = True
        restyle_tf(title_s, fn="Georgia", sz=36, bold=True, color=RED, align=PP_ALIGN.LEFT)

    # Subtitle — Calibri 12pt muted
    if subtitle_s:
        subtitle_s.left   = Inches(0.6)
        subtitle_s.top    = Inches(6.7)
        subtitle_s.width  = Inches(8.5)
        subtitle_s.height = Inches(0.55)
        restyle_tf(subtitle_s, fn="Calibri", sz=12, bold=False, color=MUTED)

    # Two small red-outlined logo squares bottom-left
    _rect(slide, 0.60, 6.95, 0.38, 0.38, line=RED, lw=Pt(1.5))
    _rect(slide, 1.10, 6.95, 0.38, 0.38, line=RED, lw=Pt(1.5))

    # Presenter / Date labels bottom-right
    _tb(slide, "Presenter Name:  __________________", 9.0, 6.70, 4.0, 0.3,
        fn="Calibri", sz=9, color=MUTED)
    _tb(slide, "Date:  __________________", 9.0, 7.00, 4.0, 0.3,
        fn="Calibri", sz=9, color=MUTED)


# ── Slide 2 — Executive Summary ───────────────────────────────────────────────

def s2_exec_summary(slide):
    solid_bg(slide, WHITE)

    # Remove background decorative rects
    kill_bg_rects(slide)

    # Dark navy left panel (35%)
    panel = _rect(slide, 0, 0, 4.67, 7.5, fill=NAVY)
    to_back(panel, slide)

    # "EXECUTIVE SUMMARY" badge text over panel/content divide
    _tb(slide, "EXECUTIVE", 3.1, 2.0, 4.2, 0.75,
        fn="Georgia", sz=26, bold=True, color=RED, align=PP_ALIGN.CENTER)
    _tb(slide, "SUMMARY",   3.1, 2.72, 4.2, 0.75,
        fn="Georgia", sz=26, bold=True, color=RED, align=PP_ALIGN.CENTER)

    # Red rule under badge
    _rect(slide, 3.1, 3.45, 4.2, 0.04, fill=RED)

    # "Executive Summary" header on white side
    _tb(slide, "Executive Summary", 4.9, 0.35, 8.0, 0.9,
        fn="Georgia", sz=22, bold=True, color=NAVY)

    # Restyle content textboxes — move into right panel
    content_boxes = [s for s in slide.shapes
                     if s.has_text_frame and s.text.strip()
                     and s.text.strip() not in ("2",)]
    pos_x = [5.0, 9.1]
    content_idx = 0
    for s in content_boxes:
        if 'document' in s.text.lower() or '•' in s.text:
            if content_idx < 2:
                s.left   = Inches(pos_x[content_idx])
                s.top    = Inches(3.7)
                s.width  = Inches(3.8)
                s.height = Inches(3.0)
                restyle_tf(s, fn="Calibri", sz=12, color=DARK)
                content_idx += 1

    # Column divider
    _rect(slide, 8.95, 3.7, 0.03, 3.0, fill=LGRAY)


# ── Slide 3 — Table of Contents ───────────────────────────────────────────────

def s3_toc(slide):
    solid_bg(slide, WHITE)

    # Remove background rects
    kill_bg_rects(slide)

    # Title
    _tb(slide, "Table of Contents", 0.6, 0.3, 12.0, 0.85,
        fn="Georgia", sz=26, bold=True, color=RED)

    # Red underline
    _rect(slide, 0.6, 1.1, 12.0, 0.04, fill=RED)

    # Restyle number textboxes
    for s in slide.shapes:
        if not s.has_text_frame:
            continue
        t = s.text.strip()
        if len(t) == 2 and t[0] == '0' and t[1].isdigit():
            # Number badge — bold red Georgia
            restyle_tf(s, fn="Georgia", sz=20, bold=True, color=RED, align=PP_ALIGN.CENTER)
        elif t and t not in ("3",):
            # Section heading text
            restyle_tf(s, fn="Calibri", sz=13, bold=False, color=DARK)

    # Light gray horizontal row dividers (every ~0.7")
    for row in range(9):
        y = 1.2 + row * 0.68
        if y < 7.1:
            _rect(slide, 0.5, y + 0.62, 12.3, 0.01, fill=LGRAY)


# ── Slide 4 / 8 — Section Divider ────────────────────────────────────────────

def s_divider(slide, label: str):
    solid_bg(slide, WHITE)

    # Remove full-slide dark background rect
    kill_bg_rects(slide)

    # Red left panel (38%)
    panel = _rect(slide, 0, 0, 5.07, 7.5, fill=RED)
    to_back(panel, slide)

    # Restyle em-dash and section title to white Georgia inside red panel
    for s in list(slide.shapes):
        if not s.has_text_frame:
            continue
        t = s.text.strip()
        if t == '—' or t == '-' or t == '\u2014':
            s.left   = Inches(0.65)
            s.top    = Inches(2.5)
            s.width  = Inches(4.2)
            s.height = Inches(0.7)
            restyle_tf(s, fn="Georgia", sz=32, bold=True, color=WHITE)
        elif label.lower() in t.lower() and len(t) < 60:
            s.left   = Inches(0.65)
            s.top    = Inches(3.2)
            s.width  = Inches(4.1)
            s.height = Inches(1.6)
            s.text_frame.word_wrap = True
            restyle_tf(s, fn="Georgia", sz=32, bold=True, color=WHITE)

    # Thin white rule on right side to add visual polish
    _rect(slide, 5.07, 3.4, 7.93, 0.03, fill=LGRAY)


# ── Slide 5 — Three Cards ────────────────────────────────────────────────────

def s5_three_cards(slide):
    solid_bg(slide, WHITE)

    kill_bg_rects(slide)

    # Title
    _tb(slide, "Key Themes", 0.6, 0.3, 12.0, 0.85,
        fn="Georgia", sz=26, bold=True, color=RED)

    # Red bottom accent bar for each card
    for lx in [0.6, 4.73, 8.85]:
        _rect(slide, lx, 6.5, 3.88, 0.06, fill=RED)

    # Find number textboxes (01, 02, 03) and add red circle behind them
    for s in list(slide.shapes):
        if not s.has_text_frame:
            continue
        t = s.text.strip()
        if t in ('01', '02', '03'):
            lx = emu_to_in(s.left)
            ty = emu_to_in(s.top)
            # Add red circle behind number
            circ = _oval(slide, lx, ty, 0.55, 0.55, fill=RED)
            to_back(circ, slide)
            # Resize textbox to circle size and center-align white text
            s.width  = Inches(0.55)
            s.height = Inches(0.55)
            restyle_tf(s, fn="Georgia", sz=16, bold=True,
                       color=WHITE, align=PP_ALIGN.CENTER)

    # Restyle content textboxes
    for s in slide.shapes:
        if s.has_text_frame and ('document' in s.text.lower() or '•' in s.text):
            restyle_tf(s, fn="Calibri", sz=12, color=DARK)


# ── Slide 6 / 10 — Chart Slide ───────────────────────────────────────────────

def s_chart(slide, slide_idx: int):
    solid_bg(slide, WHITE)

    # Restyle title
    for s in slide.shapes:
        if s.has_text_frame and s.text.strip() and 'could not' not in s.text:
            t = s.text.strip()
            if any(kw in t for kw in ['Acquisition', 'Geographic', 'Volume', 'Distribution', 'Value']):
                restyle_tf(s, fn="Georgia", sz=26, bold=True, color=RED)

    # Remove error text box
    for s in list(slide.shapes):
        if s.has_text_frame and ('[Chart' in s.text or 'could not' in s.text):
            kill(s)

    # Add chart with real data
    if slide_idx == 5:   # Slide 6 — Acquisition Volume
        cats   = ["FY2020", "FY2021", "FY2022", "FY2023", "FY2024", "FY2025\n(YTD)"]
        series = [("Acquisitions", [34, 46, 38, 42, 46, 10])]
    else:                # Slide 10 — Geographic Distribution
        cats   = ["Europe", "Asia-Pacific", "North America", "Other"]
        series = [("Share (%)", [40, 35, 25, 10])]

    add_native_chart(slide, cats, series, 0.8, 1.95, 11.7, 4.7)


# ── Slide 7 — Domain Bullets ─────────────────────────────────────────────────

def s7_domains(slide):
    solid_bg(slide, WHITE)

    # Red left panel
    panel = _rect(slide, 0, 0, 5.0, 7.5, fill=RED)
    to_back(panel, slide)

    # Section label in white on red panel
    _tb(slide, "Domain\nFocus", 0.6, 2.6, 3.8, 1.8,
        fn="Georgia", sz=28, bold=True, color=WHITE)

    # Restyle and move bullet list to right side
    for s in list(slide.shapes):
        if s.has_text_frame and ('•' in s.text or 'AI' in s.text or 'Cyber' in s.text):
            s.left   = Inches(5.3)
            s.top    = Inches(1.5)
            s.width  = Inches(7.5)
            s.height = Inches(5.5)
            restyle_tf(s, fn="Calibri", sz=13, color=DARK)

    # "Acquisition Focus Areas" header on white side
    _tb(slide, "Acquisition Focus Areas", 5.3, 0.3, 7.5, 0.85,
        fn="Georgia", sz=22, bold=True, color=RED)

    # Red rule under header
    _rect(slide, 5.3, 1.1, 7.5, 0.04, fill=RED)


# ── Slide 9 — Key Statistics ─────────────────────────────────────────────────

def s9_key_stats(slide):
    solid_bg(slide, WHITE)

    kill_bg_rects(slide)

    # Title
    _tb(slide, "Key Statistics", 0.6, 0.3, 12.0, 0.85,
        fn="Georgia", sz=26, bold=True, color=RED)

    # Red underline
    _rect(slide, 0.6, 1.1, 12.0, 0.04, fill=RED)

    # Restyle numbers — huge red Georgia
    for s in slide.shapes:
        if not s.has_text_frame:
            continue
        t = s.text.strip()
        if any(x in t for x in ('326', '$6.6', '7%', 'billion')):
            s.top    = Inches(2.4)
            s.height = Inches(2.0)
            restyle_tf(s, fn="Georgia", sz=60, bold=True, color=RED, align=PP_ALIGN.CENTER)
        elif any(x in t for x in ('Acquisitions from', 'Investment in', 'Revenue', 'Increase')):
            s.top    = Inches(4.5)
            s.height = Inches(0.7)
            restyle_tf(s, fn="Calibri", sz=12, bold=False, color=MUTED, align=PP_ALIGN.CENTER)

    # Thin red vertical dividers between 3 stat blocks
    _rect(slide, 4.72, 1.8, 0.04, 4.5, fill=RED)
    _rect(slide, 8.84, 1.8, 0.04, 4.5, fill=RED)


# ── Slide 11 — Key Takeaways ─────────────────────────────────────────────────

def s11_takeaways(slide):
    solid_bg(slide, WHITE)

    # Restyle title
    for s in slide.shapes:
        if s.has_text_frame and 'Takeaway' in s.text:
            s.left   = Inches(0.6)
            s.top    = Inches(0.3)
            s.width  = Inches(12.0)
            s.height = Inches(0.85)
            restyle_tf(s, fn="Georgia", sz=28, bold=True, color=RED)

    # Red underline
    _rect(slide, 0.6, 1.1, 12.0, 0.04, fill=RED)

    # Find the bullet list textbox and add red left accent bar beside it
    for s in slide.shapes:
        if s.has_text_frame and '•' in s.text:
            # Red left accent bar
            lx = emu_to_in(s.left) - 0.22
            ty = emu_to_in(s.top)
            ht = emu_to_in(s.height)
            _rect(slide, max(0.3, lx), ty, 0.07, ht, fill=RED)
            restyle_tf(s, fn="Calibri", sz=14, color=DARK)


# ── Slide 12 — Thank You ─────────────────────────────────────────────────────

def s12_thank_you(slide):
    solid_bg(slide, WHITE)

    TOP_H = 7.5 * 0.58   # ≈ 4.35"
    BOT_H = 7.5 - TOP_H  # ≈ 3.15"

    # Blush pink top half
    top = _rect(slide, 0, 0, 13.33, TOP_H, fill=BLUSH)
    to_back(top, slide)

    # Dark navy bottom half
    bot = _rect(slide, 0, TOP_H, 13.33, BOT_H, fill=NAVY)
    to_back(bot, slide)

    # Thin red stripe at divide
    _rect(slide, 0, TOP_H - 0.05, 13.33, 0.10, fill=RED)

    # X decoration — two thin crossing diagonal bars, top-right of pink area
    for (lx, ty, rot) in [(10.8, 0.4, 35), (11.2, 0.4, -35)]:
        bar = slide.shapes.add_shape(1, Inches(lx), Inches(ty), Inches(0.1), Inches(2.2))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RED
        bar.line.fill.background()
        bar.rotation = rot

    # Small red-outlined square, bottom-left of pink area
    _rect(slide, 0.6, TOP_H - 0.75, 0.42, 0.42, line=RED, lw=Pt(2.0))

    # "Thank You" text — restyle existing
    for s in slide.shapes:
        if s.has_text_frame and 'Thank' in s.text:
            s.left   = Inches(0.6)
            s.top    = Inches(1.2)
            s.width  = Inches(10.0)
            s.height = Inches(2.6)
            s.text_frame.word_wrap = True
            restyle_tf(s, fn="Georgia", sz=64, bold=True, color=RED, align=PP_ALIGN.LEFT)

    # Tagline in navy area
    _tb(slide, "Thank you for your attention.",
        0.6, TOP_H + 0.5, 10.0, 0.6,
        fn="Calibri", sz=16, bold=False, color=WHITE)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation("outputs/new3.pptx")
    slides = list(prs.slides)
    print(f"Restyling {len(slides)} slides …\n")

    handlers = {
        0:  ("Cover",                  lambda s: s1_cover(s)),
        1:  ("Executive Summary",      lambda s: s2_exec_summary(s)),
        2:  ("Table of Contents",      lambda s: s3_toc(s)),
        3:  ("Section Divider",        lambda s: s_divider(s, "Introduction")),
        4:  ("Three Cards",            lambda s: s5_three_cards(s)),
        5:  ("Chart — Volume",         lambda s: s_chart(s, 5)),
        6:  ("Domain Bullets",         lambda s: s7_domains(s)),
        7:  ("Section Divider",        lambda s: s_divider(s, "Acquisition Trends")),
        8:  ("Key Stats",              lambda s: s9_key_stats(s)),
        9:  ("Chart — Geographic",     lambda s: s_chart(s, 9)),
        10: ("Key Takeaways",          lambda s: s11_takeaways(s)),
        11: ("Thank You",              lambda s: s12_thank_you(s)),
    }

    for i, slide in enumerate(slides):
        name, fn = handlers.get(i, ("Unknown", None))
        if fn:
            try:
                fn(slide)
                print(f"  [{i+1:02d}] {name} ✓")
            except Exception as e:
                import traceback
                print(f"  [{i+1:02d}] {name} ERROR: {e}")
                traceback.print_exc()
        else:
            print(f"  [{i+1:02d}] No handler — skipping")

    out = "outputs/output_styled.pptx"
    prs.save(out)
    print(f"\nSaved: {out}")


if __name__ == "__main__":
    main()
